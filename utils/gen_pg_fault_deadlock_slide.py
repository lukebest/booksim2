#!/usr/bin/env python3
"""One-slide summary: how the 8x6 PG NoC study tolerates faults and breaks deadlock.

Layout is described once as primitive ops, then rendered by two backends:
  --pptx  editable 16:9 slide (needs python-pptx)
  --png   preview raster (needs matplotlib)

Run the pptx backend from the venv that has python-pptx:
  .venv-ppt/bin/python utils/gen_pg_fault_deadlock_slide.py --pptx
  python3 utils/gen_pg_fault_deadlock_slide.py --png
"""
from __future__ import annotations

import argparse
import textwrap
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
OUT_PPTX = ROOT / "results" / "pg_fault_deadlock_slide.pptx"
OUT_PNG = ROOT / "results" / "pg_fault_deadlock_slide.png"

SLIDE_W, SLIDE_H = 13.333, 7.5
FONT = "Microsoft YaHei"

INK = "1F2933"
RED = "C7000B"
BLUE = "0B6BA8"
GREEN = "1E7F4F"
ORANGE = "C87209"
GREY = "5A6673"
GREY_L = "9AA5B1"
WHITE = "FFFFFF"
CARD_BG = "F7F8FA"
CARD_LN = "D9DEE4"
HDR_BG = "E8EDF2"


class Deck:
    """Collects resolution-independent drawing ops in inches (y grows down)."""

    def __init__(self) -> None:
        self.ops: list[dict] = []

    def rect(self, x, y, w, h, *, fill=None, line=None, lw=0.75, round_=0.0):
        self.ops.append(dict(kind="rect", x=x, y=y, w=w, h=h, fill=fill,
                             line=line, lw=lw, round_=round_))

    def oval(self, cx, cy, r, *, fill=None, line=None, lw=0.75):
        self.ops.append(dict(kind="oval", cx=cx, cy=cy, r=r, fill=fill,
                             line=line, lw=lw))

    def star(self, cx, cy, r, *, fill=RED, line=None, lw=0.75):
        self.ops.append(dict(kind="star", cx=cx, cy=cy, r=r, fill=fill,
                             line=line, lw=lw))

    def line(self, x1, y1, x2, y2, *, color=GREY, lw=1.0, dash=False,
             arrow=False):
        self.ops.append(dict(kind="line", x1=x1, y1=y1, x2=x2, y2=y2,
                             color=color, lw=lw, dash=dash, arrow=arrow))

    def cross(self, cx, cy, r, *, color=RED, lw=1.6):
        self.line(cx - r, cy - r, cx + r, cy + r, color=color, lw=lw)
        self.line(cx - r, cy + r, cx + r, cy - r, color=color, lw=lw)

    def text(self, x, y, w, h, paras, *, valign="top"):
        self.ops.append(dict(kind="text", x=x, y=y, w=w, h=h,
                             paras=paras, valign=valign))


def p(t, *, size=11.0, bold=False, color=INK, align="l", space=3.0,
      spacing=1.0):
    return dict(t=t, size=size, bold=bold, color=color, align=align,
                space=space, spacing=spacing)


def card(d: Deck, x, y, w, h, title, *, accent=BLUE, hdr_h=0.34):
    d.rect(x, y, w, h, fill=CARD_BG, line=CARD_LN, lw=0.75, round_=0.04)
    d.rect(x, y, w, hdr_h, fill=HDR_BG, line=None, round_=0.04)
    d.rect(x, y, 0.055, hdr_h, fill=accent, line=None)
    d.text(x + 0.16, y + 0.03, w - 0.3, hdr_h,
           [p(title, size=12.5, bold=True, color=INK, space=0)],
           valign="middle")


# --------------------------------------------------------------------------
# Figure 1: fault tolerance on the residual graph
# --------------------------------------------------------------------------

def fig_fault(d: Deck, x0, y0, w, h):
    nx, ny = 8, 6
    sx = w / (nx - 1 + 1.2)
    sy = h / (ny - 1 + 1.6)
    ox = x0 + sx * 0.6
    oy = y0 + sy * 0.5
    r = min(sx, sy) * 0.20

    def pos(cx, cy):
        return ox + cx * sx, oy + cy * sy

    dead = {(3, 2), (4, 2)}
    dead_link = ((1, 4), (2, 4))
    sac = (3, 3)
    path = [(1, 1), (2, 1), (3, 1), (4, 1), (5, 1),
            (5, 2), (5, 3), (6, 3), (6, 4)]

    for cy in range(ny):
        for cx in range(nx):
            if (cx, cy) in dead:
                continue
            for dx, dy in ((1, 0), (0, 1)):
                bx, by = cx + dx, cy + dy
                if bx >= nx or by >= ny or (bx, by) in dead:
                    continue
                if {(cx, cy), (bx, by)} == {*dead_link}:
                    continue
                ax_, ay_ = pos(cx, cy)
                bx_, by_ = pos(bx, by)
                d.line(ax_, ay_, bx_, by_, color="CBD2D9", lw=0.7)

    ax_, ay_ = pos(*dead_link[0])
    bx_, by_ = pos(*dead_link[1])
    d.line(ax_, ay_, bx_, by_, color=RED, lw=1.1, dash=True)
    d.cross((ax_ + bx_) / 2, (ay_ + by_) / 2, r * 0.55, color=RED, lw=1.2)

    for i in range(len(path) - 1):
        ax_, ay_ = pos(*path[i])
        bx_, by_ = pos(*path[i + 1])
        d.line(ax_, ay_, bx_, by_, color=GREEN, lw=2.1,
               arrow=(i == len(path) - 2))

    for cy in range(ny):
        for cx in range(nx):
            px, py = pos(cx, cy)
            if (cx, cy) in dead:
                d.oval(px, py, r, fill=WHITE, line=RED, lw=1.1)
                d.cross(px, py, r * 0.62, color=RED, lw=1.3)
            elif (cx, cy) == sac:
                d.oval(px, py, r * 1.15, fill=WHITE, line=ORANGE, lw=1.5)
            elif (cx, cy) in (path[0], path[-1]):
                d.oval(px, py, r * 1.1, fill=GREEN, line=GREEN, lw=0.8)
            else:
                d.oval(px, py, r * 0.8, fill="AEB8C2", line=None, lw=0)

    sx_, sy_ = pos(*path[0])
    dx_, dy_ = pos(*path[-1])
    d.text(sx_ - 0.30, sy_ - 0.34, 0.6, 0.2,
           [p("S", size=8.5, bold=True, color=GREEN, align="c", space=0)])
    d.text(dx_ - 0.30, dy_ + 0.10, 0.6, 0.2,
           [p("D", size=8.5, bold=True, color=GREEN, align="c", space=0)])
    d.text(x0, y0 + h - 0.20, w, 0.22,
           [p("× 故障　○ 牺牲　→ 残图上的绕行路径",
              size=8.2, color=GREY, align="c", space=0)])


# --------------------------------------------------------------------------
# Figure 2: channel dependency graph cycle
# --------------------------------------------------------------------------

def fig_cdg(d: Deck, x0, y0, w, h):
    bw, bh = w * 0.30, 0.30
    cx0, cy0 = x0 + w * 0.06, y0 + 0.06
    cx1 = x0 + w - bw - w * 0.06
    cy1 = y0 + h - bh - 0.34

    boxes = {
        "c1": (cx0, cy0),
        "c2": (cx1, cy0),
        "c3": (cx1, cy1),
        "c4": (cx0, cy1),
    }
    for name, (bx, by) in boxes.items():
        d.rect(bx, by, bw, bh, fill=WHITE, line=BLUE, lw=1.0, round_=0.05)
        d.text(bx, by + 0.03, bw, bh,
               [p(name, size=9.5, bold=True, color=BLUE, align="c", space=0)],
               valign="middle")

    def mid(name, side):
        bx, by = boxes[name]
        if side == "r":
            return bx + bw, by + bh / 2
        if side == "l":
            return bx, by + bh / 2
        if side == "b":
            return bx + bw / 2, by + bh
        return bx + bw / 2, by

    for a, sa, b, sb in (("c1", "r", "c2", "l"),
                         ("c2", "b", "c3", "t"),
                         ("c3", "l", "c4", "r")):
        x1, y1 = mid(a, sa)
        x2, y2 = mid(b, sb)
        d.line(x1, y1, x2, y2, color=RED, lw=1.6, arrow=True)

    x1, y1 = mid("c4", "t")
    x2, y2 = mid("c1", "b")
    d.line(x1, y1, x2, y2, color=GREY_L, lw=1.4, dash=True, arrow=True)
    d.cross(x1, (y1 + y2) / 2, 0.10, color=RED, lw=1.9)
    d.text(x1 + 0.16, (y1 + y2) / 2 - 0.20, w * 0.60, 0.40,
           [p("A / B / C", size=8.2, bold=True, color=RED, space=1.0),
            p("断掉这一处依赖", size=8.2, bold=True, color=RED, space=0)])

    d.text(x0, y0 + h - 0.24, w, 0.24,
           [p("有环 ⇒ 死锁　|　无环 ⇒ 安全",
              size=8.8, bold=True, color=INK, align="c", space=0)])


# --------------------------------------------------------------------------
# Mini figures for the three deadlock-freedom constructions
# --------------------------------------------------------------------------

def fig_turn(d: Deck, x0, y0, w, h):
    cx, cy = x0 + w * 0.42, y0 + h * 0.52
    d.oval(cx, cy, 0.075, fill=BLUE, line=None)
    d.line(cx, y0 + h * 0.94, cx, cy + 0.10, color=BLUE, lw=1.6, arrow=True)
    d.line(cx + 0.10, cy, x0 + w * 0.95, cy, color=GREEN, lw=1.6, arrow=True)
    d.line(cx, cy - 0.10, cx, y0 + h * 0.12, color=GREY_L, lw=1.5,
           dash=True, arrow=True)
    d.cross(cx, y0 + h * 0.26, 0.075, color=RED, lw=1.7)
    d.text(cx - 0.46, y0 + h * 0.26 - 0.10, 0.38, 0.2,
           [p("禁", size=8.2, bold=True, color=RED, align="r", space=0)])


def fig_vc(d: Deck, x0, y0, w, h):
    bh = h * 0.24
    for i, (lab, col, yy) in enumerate((
            ("VC1", ORANGE, y0 + h * 0.16),
            ("VC0", BLUE, y0 + h * 0.60))):
        d.rect(x0 + w * 0.14, yy, w * 0.80, bh, fill=WHITE, line=col, lw=1.0,
               round_=0.03)
        d.line(x0 + w * 0.22, yy + bh / 2, x0 + w * 0.86, yy + bh / 2,
               color=col, lw=1.4, arrow=True)
        d.text(x0, yy, w * 0.13, bh,
               [p(lab, size=7.5, bold=True, color=col, align="r", space=0)],
               valign="middle")
    d.line(x0 + w * 0.52, y0 + h * 0.60, x0 + w * 0.52, y0 + h * 0.16 + bh,
           color=GREEN, lw=1.5, arrow=True)
    d.text(x0 + w * 0.56, y0 + h * 0.44, w * 0.45, 0.2,
           [p("只升不降", size=7.2, bold=True, color=GREEN, space=0)])


def fig_batch(d: Deck, x0, y0, w, h):
    ytop = y0 + h * 0.30
    bh = h * 0.26
    w1 = w * 0.34
    w2 = w * 0.30
    d.rect(x0 + w * 0.06, ytop, w1, bh, fill=WHITE, line=BLUE, lw=1.0,
           round_=0.03)
    d.text(x0 + w * 0.06, ytop, w1, bh,
           [p("批 1", size=7.8, bold=True, color=BLUE, align="c", space=0)],
           valign="middle")
    bx = x0 + w * 0.06 + w1 + w * 0.12
    d.rect(bx, ytop, w2, bh, fill=WHITE, line=BLUE, lw=1.0, round_=0.03)
    d.text(bx, ytop, w2, bh,
           [p("批 2", size=7.8, bold=True, color=BLUE, align="c", space=0)],
           valign="middle")
    bar = x0 + w * 0.06 + w1 + w * 0.06
    d.line(bar, ytop - 0.06, bar, ytop + bh + 0.06, color=RED, lw=1.8)
    d.text(bar - w * 0.30, ytop + bh + 0.06, w * 0.60, 0.2,
           [p("barrier", size=7.5, bold=True, color=RED, align="c", space=0)])
    d.line(x0 + w * 0.06, ytop + bh + 0.30, x0 + w * 0.94, ytop + bh + 0.30,
           color=GREY_L, lw=0.9, arrow=True)
    d.text(x0, ytop + bh + 0.32, w, 0.2,
           [p("时间", size=7.5, color=GREY, align="c", space=0)])


# --------------------------------------------------------------------------
# Slide composition
# --------------------------------------------------------------------------

def build() -> Deck:
    d = Deck()
    d.rect(0, 0, SLIDE_W, SLIDE_H, fill=WHITE, line=None)

    # Title band
    d.rect(0, 0, SLIDE_W, 0.92, fill=INK, line=None)
    d.rect(0, 0, 0.10, 0.92, fill=RED, line=None)
    d.text(0.34, 0.10, 10.6, 0.44,
           [p("8×6 NoC 容错路由：如何容错 + 如何解死锁",
              size=21, bold=True, color=WHITE, space=0)])
    d.text(0.36, 0.55, 12.6, 0.30,
           [p("预算故障目录 ≤4 router + ≤8 link（44 场景）· dispatch alltoall · "
              "三条硬性质：避障可达 / 无死锁 / 保序",
              size=10.0, color="C3CBD4", space=0)])

    # Row 1 -----------------------------------------------------------------
    ry, rh = 1.06, 3.32
    cw1 = 6.28
    card(d, 0.30, ry, cw1, rh, "① 如何容错：残图重算路由，兜底才牺牲", accent=GREEN)
    fig_fault(d, 0.44, ry + 0.40, 2.72, 2.62)
    d.text(3.28, ry + 0.42, 3.16, rh - 0.52, [
        p("故障模型：dead＝router 与链路同死；transit＝PE 退出计算、"
          "router 仍替别人转发。", size=9.0, color=INK, space=4.0),
        p("① 删点删边 → 残图", size=9.6, bold=True, color=GREEN, space=1.0),
        p("故障不进邻接表，约束 BFS / Dijkstra 只在存活图上搜，路径不穿洞。",
          size=8.8, color=GREY, space=4.0),
        p("② 整表校验", size=9.6, bold=True, color=GREEN, space=1.0),
        p("每对唯一 · 每跳合法 · CDG 无环，三条全过才可用。",
          size=8.8, color=GREY, space=4.0),
        p("③ 建不出表才牺牲", size=9.6, bold=True, color=GREEN, space=1.0),
        p("最小基数逐级放宽：孤立点 → 单点 → 点对 → 整行列 → 健康矩形。",
          size=8.8, color=GREY, space=4.0),
        p("④ 牺牲不免费", size=9.6, bold=True, color=ORANGE, space=1.0),
        p("A 变小后按 m_eff = ⌈m0·(48/A)²⌉ 重标定，强扩展惩罚重牺牲。",
          size=8.8, color=ORANGE, space=0),
    ])

    cx2 = 0.30 + cw1 + 0.16
    cw2 = SLIDE_W - cx2 - 0.30
    card(d, cx2, ry, cw2, rh, "② 如何解死锁：让通道依赖图（CDG）无环", accent=RED)
    fig_cdg(d, cx2 + 0.16, ry + 0.40, 2.62, 2.62)
    d.text(cx2 + 2.96, ry + 0.42, cw2 - 3.12, rh - 0.52, [
        p("死锁 = 几条路径各占一段、又互等对方的下一段，缓冲有限即永久卡住。",
          size=9.0, color=INK, space=4.0),
        p("判据（Dally–Seitz）", size=9.6, bold=True, color=RED, space=1.0),
        p("以（有向通道, VC）为节点、路径相邻两跳为边建出 CDG；"
          "无环 ⇒ 不存在循环等待 ⇒ 任意注入都不死锁。",
          size=8.8, color=GREY, space=4.0),
        p("三重保证", size=9.6, bold=True, color=RED, space=1.0),
        p("① 构造时就断掉一类依赖（下方 A / B / C）；"
          "② 统一用 build_cdg + cdg_acyclic 事后校验；"
          "③ 负载均衡换路后整表重验，出环立即回退。",
          size=8.8, color=GREY, space=4.0),
        p("→ 44 场景 × 全部方案：CDG 成环 0 例；M3 Up*/Down* 在 36/36 "
          "目录格零牺牲可行。", size=8.8, bold=True, color=GREEN, space=0),
    ])

    # Row 2 -----------------------------------------------------------------
    ty, th = 4.52, 1.72
    tw = (SLIDE_W - 0.60 - 0.32) / 3
    specs = [
        ("A 转向禁令 · 1 VC", BLUE, fig_turn, [
            p("永久禁掉「会闭合环」的那类转弯，不花 VC。",
              size=8.6, color=GREY, space=3.0),
            p("Up*/Down*：树高度标号，禁 down→up；XY：禁 Y→X；"
              "Glass–Ni：每层一个最小转向集。",
              size=8.6, color=INK, space=3.0),
            p("代价：绕行方向变少，大洞时可能要牺牲。",
              size=8.2, color=ORANGE, space=0),
        ]),
        ("B VC 分层 · 2–3 VC", ORANGE, fig_vc, [
            p("把会互锁的路径分到不同虚通道，层内各自无环。",
              size=8.6, color=GREY, space=3.0),
            p("LASH：贪心装入最少无环层；Dual UD：VC0 走 UD、VC1 走 DU。",
              size=8.6, color=INK, space=3.0),
            p("整条路径锁一层、层号只升不降 ⇒ 并集无环。",
              size=8.2, color=GREEN, space=0),
        ]),
        ("C 时间分批 + barrier · 1 VC", GREEN, fig_batch, [
            p("同一划分用时间代替 VC：按 OD 整层切开，逐批串行。",
              size=8.6, color=INK, space=3.0),
            p("批间取残图图中心 gather→broadcast，"
              "T_sync = 2·radius_wire（H=7 / V=9，同仿真口径）。",
              size=8.6, color=GREY, space=3.0),
            p("批内 CDG 无环 ⇒ 面积按 1 VC 计。",
              size=8.2, color=GREEN, space=0),
        ]),
    ]
    for i, (title, accent, fig, paras) in enumerate(specs):
        cx = 0.30 + i * (tw + 0.16)
        card(d, cx, ty, tw, th, title, accent=accent, hdr_h=0.32)
        fig(d, cx + 0.10, ty + 0.38, 1.02, th - 0.48)
        d.text(cx + 1.20, ty + 0.38, tw - 1.32, th - 0.44, paras)

    # Bottom band -----------------------------------------------------------
    by, bh = 6.40, 0.90
    d.rect(0.30, by, SLIDE_W - 0.60, bh, fill="FBF0F0", line=RED, lw=0.9,
           round_=0.04)
    d.rect(0.30, by, 0.055, bh, fill=RED, line=None)
    d.text(0.48, by + 0.10, 5.80, bh - 0.16, [
        p("第三条性质：保序", size=9.6, bold=True, color=RED, space=1.5),
        p("每对源宿只有一条离线定死的路径，运行时不自适应、不中途换 VC，"
          "flit 顺序天然保持；分批方案 = 批内唯一路径 + 批间屏障。",
          size=8.8, color=INK, space=0),
    ])
    d.text(6.55, by + 0.10, SLIDE_W - 6.55 - 0.45, bh - 0.16, [
        p("端到端结论（44 场景最差值，已含牺牲后的强扩展重标定）",
          size=9.6, bold=True, color=RED, space=1.5),
        p("轻载 m0=1：M3′ Up*/Down* 790 ns（1 VC，area 0.90）→ Super-turn "
          "635 ns（2 VC，1.24）；重载 m0=13：分批屏障 BB UD policy 7061 ns"
          "（1 VC）最优。", size=8.8, color=INK, space=0),
    ])

    return d


# --------------------------------------------------------------------------
# Backends
# --------------------------------------------------------------------------

def emit_pptx(deck: Deck, path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches, Pt

    ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
    ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE}

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    def style(shape, fill, line, lw):
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill)
        else:
            shape.fill.background()
        if line:
            shape.line.color.rgb = RGBColor.from_string(line)
            shape.line.width = Pt(lw)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False

    for op in deck.ops:
        k = op["kind"]
        if k == "rect":
            kind = (MSO_SHAPE.ROUNDED_RECTANGLE if op["round_"]
                    else MSO_SHAPE.RECTANGLE)
            sh = shapes.add_shape(kind, Inches(op["x"]), Inches(op["y"]),
                                  Inches(op["w"]), Inches(op["h"]))
            if op["round_"]:
                adj = op["round_"] / max(min(op["w"], op["h"]), 1e-6)
                sh.adjustments[0] = min(max(adj, 0.0), 0.5)
            style(sh, op["fill"], op["line"], op["lw"])
            sh.text_frame.word_wrap = False
        elif k in ("oval", "star"):
            kind = (MSO_SHAPE.OVAL if k == "oval"
                    else MSO_SHAPE.STAR_5_POINT)
            sh = shapes.add_shape(
                kind, Inches(op["cx"] - op["r"]),
                Inches(op["cy"] - op["r"]), Inches(op["r"] * 2),
                Inches(op["r"] * 2))
            style(sh, op["fill"], op["line"], op["lw"])
        elif k == "line":
            conn = shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(op["x1"]), Inches(op["y1"]),
                Inches(op["x2"]), Inches(op["y2"]))
            conn.line.color.rgb = RGBColor.from_string(op["color"])
            conn.line.width = Pt(op["lw"])
            ln = conn.line._get_or_add_ln()
            if op["dash"]:
                dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
                ln.append(dash)
            if op["arrow"]:
                tail = ln.makeelement(qn("a:tailEnd"),
                                      {"type": "triangle", "w": "med",
                                       "len": "med"})
                ln.append(tail)
        elif k == "text":
            box = shapes.add_textbox(Inches(op["x"]), Inches(op["y"]),
                                     Inches(op["w"]), Inches(op["h"]))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = ANCHOR[op["valign"]]
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            for i, pa in enumerate(op["paras"]):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.alignment = ALIGN[pa["align"]]
                para.space_after = Pt(pa["space"])
                para.line_spacing = pa["spacing"] * 1.18
                run = para.add_run()
                run.text = pa["t"]
                run.font.size = Pt(pa["size"])
                run.font.bold = pa["bold"]
                run.font.name = FONT
                run.font.color.rgb = RGBColor.from_string(pa["color"])

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print(f"wrote {path}")


def emit_png(deck: Deck, path: Path, dpi: int = 170) -> None:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import Circle, FancyBboxPatch, Rectangle

    plt.rcParams["font.sans-serif"] = ["WenQuanYi Micro Hei",
                                       "Droid Sans Fallback", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False

    fig = plt.figure(figsize=(SLIDE_W, SLIDE_H))
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, SLIDE_W)
    ax.set_ylim(SLIDE_H, 0)
    ax.axis("off")

    def c(h):
        return f"#{h}"

    for op in deck.ops:
        k = op["kind"]
        if k == "rect":
            fc = c(op["fill"]) if op["fill"] else "none"
            ec = c(op["line"]) if op["line"] else "none"
            if op["round_"]:
                pad = op["round_"]
                ax.add_patch(FancyBboxPatch(
                    (op["x"] + pad, op["y"] + pad),
                    op["w"] - 2 * pad, op["h"] - 2 * pad,
                    boxstyle=f"round,pad={pad},rounding_size={pad}",
                    facecolor=fc, edgecolor=ec, linewidth=op["lw"]))
            else:
                ax.add_patch(Rectangle((op["x"], op["y"]), op["w"], op["h"],
                                       facecolor=fc, edgecolor=ec,
                                       linewidth=op["lw"]))
        elif k == "oval":
            fc = c(op["fill"]) if op["fill"] else "none"
            ec = c(op["line"]) if op["line"] else "none"
            ax.add_patch(Circle((op["cx"], op["cy"]), op["r"], facecolor=fc,
                                edgecolor=ec, linewidth=op["lw"]))
        elif k == "star":
            fc = c(op["fill"]) if op["fill"] else "none"
            ax.plot([op["cx"]], [op["cy"]], marker="*",
                    markersize=op["r"] * 2 * 72, color=fc, linestyle="none")
        elif k == "line":
            style = dict(color=c(op["color"]), lw=op["lw"],
                         linestyle=(0, (3, 2)) if op["dash"] else "-")
            if op["arrow"]:
                ax.annotate("", xy=(op["x2"], op["y2"]),
                            xytext=(op["x1"], op["y1"]),
                            arrowprops=dict(arrowstyle="-|>", shrinkA=0,
                                            shrinkB=0, mutation_scale=9,
                                            **style))
            else:
                ax.plot([op["x1"], op["x2"]], [op["y1"], op["y2"]], **style)
        elif k == "text":
            heights = []
            wrapped = []
            for pa in op["paras"]:
                cpl = max(int(op["w"] / (pa["size"] / 72.0 * 0.95)), 6)
                lines = textwrap.wrap(pa["t"], cpl) or [""]
                lh = pa["size"] / 72.0 * 1.30 * pa["spacing"]
                wrapped.append((pa, lines, lh))
                heights.append(len(lines) * lh + pa["space"] / 72.0)
            total = sum(heights)
            y = (op["y"] if op["valign"] == "top"
                 else op["y"] + (op["h"] - total) / 2)
            for pa, lines, lh in wrapped:
                for ln_txt in lines:
                    if pa["align"] == "l":
                        tx, ha = op["x"], "left"
                    elif pa["align"] == "c":
                        tx, ha = op["x"] + op["w"] / 2, "center"
                    else:
                        tx, ha = op["x"] + op["w"], "right"
                    ax.text(tx, y + lh * 0.78, ln_txt, ha=ha, va="baseline",
                            fontsize=pa["size"], color=c(pa["color"]),
                            fontweight="bold" if pa["bold"] else "normal")
                    y += lh
                y += pa["space"] / 72.0

    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, dpi=dpi)
    plt.close(fig)
    print(f"wrote {path}")


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", action="store_true")
    ap.add_argument("--png", action="store_true")
    args = ap.parse_args()
    if not (args.pptx or args.png):
        args.pptx = args.png = True

    deck = build()
    if args.pptx:
        emit_pptx(deck, OUT_PPTX)
    if args.png:
        emit_png(deck, OUT_PNG)


if __name__ == "__main__":
    main()
