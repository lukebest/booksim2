#!/usr/bin/env python3
"""Build the architecture-review deck as a real .pptx (Huawei corporate style).

Same content and same numbers as the HTML deck in ppt/, rebuilt with
python-pptx so the file can be opened and edited in PowerPoint / WPS /
Keynote. Every figure is read from ppt/images/, which is produced by
utils/ppt_ring2_figs.py from results/*.json.

Inline markup accepted by the text helpers:
    **bold**      bold run
    [[red]]       brand-red bold run

Usage:
    python3 utils/build_ring2_pptx.py            # writes ppt/ring2-write-fairness.pptx
"""
from __future__ import annotations

import json
import math
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "ppt" / "images"
OUT = ROOT / "ppt" / "ring2-write-fairness.pptx"
DECK_JSON = ROOT / "results" / "deck_ring2_data.json"

# ---------------------------------------------------------------- palette
PAPER = RGBColor(0xF4, 0xF7, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x20, 0x24, 0x2C)
MUTED = RGBColor(0x66, 0x70, 0x80)
LINE = RGBColor(0xDC, 0xE3, 0xEA)
RED = RGBColor(0xD2, 0x0A, 0x2E)
RED_DARK = RGBColor(0x9C, 0x00, 0x1F)
GREY_BG = RGBColor(0xED, 0xF2, 0xF7)
DARK = RGBColor(0x14, 0x17, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Microsoft YaHei"
FONT_EN = "Arial"

W, H = 13.333, 7.5
ML = 0.62                      # left / right page margin
CW = W - 2 * ML                # content width
TOP = 1.06                     # first usable y below the chrome rule
BOT = 6.92                     # last usable y above the page number


# ---------------------------------------------------------------- helpers
def _units(s: str) -> float:
    """Width of a string in em, counting CJK as full width."""
    return sum(1.0 if ord(c) > 0x2E80 else 0.52 for c in s)


def _plain(s: str) -> str:
    return s.replace("**", "").replace("[[", "").replace("]]", "")


# A paragraph's line_spacing multiplies the font's own line height, which for
# the CJK faces used here is ~1.46 em rather than 1.0 em. Every height estimate
# below has to carry that factor or the text silently overruns its box.
LH = 1.46


def n_lines(lines: list[str], box_w: float, pt: float) -> int:
    per_line = box_w / (pt / 72)
    return sum(max(1, math.ceil(_units(_plain(t)) / per_line)) for t in lines)


def block_h(lines: list[str], box_w: float, pt: float, lsp: float,
            gap: float) -> float:
    return (n_lines(lines, box_w, pt) * pt * lsp * LH
            + (len(lines) - 1) * gap) / 72


def fit_pt(lines: list[str], box_w: float, box_h: float, base: float,
           floor: float = 8, lsp: float = 1.0, gap: float = 5.0) -> float:
    """Largest point size at which `lines` still fit in a box_w x box_h inch box."""
    pt = float(base)
    while pt > floor:
        if block_h(lines, box_w, pt, lsp, gap) <= box_h:
            return pt
        pt -= 0.5
    return floor


_TOKEN = re.compile(r"(\*\*.+?\*\*|\[\[.+?\]\])")


def _runs(par, text: str, pt: int, color: RGBColor, bold: bool = False,
          font: str = FONT) -> None:
    for chunk in _TOKEN.split(text):
        if not chunk:
            continue
        if chunk.startswith("**"):
            body, b, c = chunk[2:-2], True, color
        elif chunk.startswith("[["):
            body, b, c = chunk[2:-2], True, RED
        else:
            body, b, c = chunk, bold, color
        r = par.add_run()
        r.text = body
        r.font.size = Pt(pt)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = font


def textbox(slide, x, y, w, h, lines, pt, color=INK, bold=False, lsp=1.0,
            gap=5.0, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, text in enumerate(lines):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = align
        par.line_spacing = lsp
        if i:
            par.space_before = Pt(gap)
        _runs(par, text, pt, color, bold, font)
    return tb


def autotext(slide, x, y, w, h, lines, base, color=INK, bold=False, floor=8,
             lsp=1.0, gap=5.0, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT):
    pt = fit_pt(lines, w, h, base, floor, lsp, gap)
    return textbox(slide, x, y, w, h, lines, pt, color, bold, lsp, gap, align,
                   anchor, font)


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0,
         shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.text_frame.word_wrap = True
    return sp


def image(slide, path: Path, x, y, w, h):
    """Place an image contained inside the x/y/w/h box, centred.

    Returns the box actually occupied, so captions can hug the figure instead
    of the (often much taller) slot reserved for it.
    """
    from PIL import Image as PILImage
    iw, ih = PILImage.open(path).size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    dx, dy = x + (w - dw) / 2, y + (h - dh) / 2
    slide.shapes.add_picture(str(path), Inches(dx), Inches(dy), Inches(dw),
                             Inches(dh))
    return dx, dy, dw, dh


# ------------------------------------------------------------ page furniture
def new_slide(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, fill=bg)
    return s


def brand(slide, on_dark=False):
    rect(slide, W - ML - 1.62, 0.30, 0.20, 0.20, fill=RED)
    textbox(slide, W - ML - 1.34, 0.27, 1.34, 0.28, ["HUAWEI"], 13,
            WHITE if on_dark else INK, bold=True, font=FONT_EN)


def chrome(slide, title: str):
    rect(slide, ML, 0.33, 0.15, 0.15, fill=RED)
    textbox(slide, ML + 0.26, 0.29, W - 2 * ML - 2.0, 0.26, [title], 11.5, MUTED,
            bold=True)
    rect(slide, ML, 0.70, CW, 0.011, fill=LINE)


def page_no(slide, i, total, on_dark=False):
    textbox(slide, ML, 7.06, 1.6, 0.24, [f"{i:02d} / {total:02d}"], 10,
            RGBColor(0x55, 0x5B, 0x66) if on_dark else RGBColor(0xA8, 0xB0, 0xBA),
            font=FONT_EN)


def kicker(slide, x, y, w, text, color=RED, pt=11):
    spaced = " ".join(text)
    return textbox(slide, x, y, w, 0.26, [spaced],
                   fit_pt([spaced], w, 0.26, pt, 7.5), color, bold=True,
                   font=FONT_EN)


def rule(slide, x, y, w=1.15):
    rect(slide, x, y, w, 0.045, fill=RED)


def band(slide, x, y, w, h, text, base=15):
    rect(slide, x, y, w, h, fill=RED)
    autotext(slide, x + 0.34, y + 0.10, w - 0.68, h - 0.20, [text], base,
             WHITE, floor=9, lsp=1.04, anchor=MSO_ANCHOR.MIDDLE)


def card_need(title, bodies, cw, tpt, bpt, pad, num=False):
    """Height a card needs to hold its text at full size, so a row of cards can
    be sized to its content instead of stretching to fill the slide."""
    h = 2 * pad + (0.58 if num else 0.0)
    if title:
        h += block_h([title], cw, tpt, 1.0, 0) + 0.13
    if bodies:
        h += block_h(bodies, cw, bpt, BODY_LSP, BODY_GAP)
    return h


def row_h(cells, cwid, tpt, bpt, pad, avail, num=False, floor_h=1.5):
    need = max(card_need(c.get("t", ""), c.get("b", []), cwid - 2 * pad, tpt,
                         bpt, pad, num) for c in cells)
    return min(avail, max(floor_h, need))


def card(slide, x, y, w, h, accent=False, flat=False):
    return rect(slide, x, y, w, h, fill=GREY_BG if flat else CARD,
                line=RED if accent else LINE, lw=1.5 if accent else 1.0)


BODY_LSP = 1.02
BODY_GAP = 6.0


def card_text(slide, x, y, w, h, title, bodies, accent=False, flat=False,
              tpt=14, bpt=11.5, pad=0.24, num=None, floor=8.5):
    card(slide, x, y, w, h, accent, flat)
    cx, cw = x + pad, w - 2 * pad
    inner = h - 2 * pad
    num_h = 0.58 if num is not None else 0.0

    th = tgap = 0.0
    tp = 0.0
    if title:
        tp = fit_pt([title], cw, 0.78, tpt, 9.5, lsp=1.0, gap=0)
        th = block_h([title], cw, tp, 1.0, 0)
        tgap = 0.13

    bh = 0.0
    bp = 0.0
    if bodies:
        bp = fit_pt(bodies, cw, inner - num_h - th - tgap, bpt, floor,
                    lsp=BODY_LSP, gap=BODY_GAP)
        bh = block_h(bodies, cw, bp, BODY_LSP, BODY_GAP)

    # Numbered cards keep their pill on a shared baseline across the row; only
    # the text below it is optically centred in whatever space is left.
    cy = y + pad
    if num is not None:
        rect(slide, cx, cy, 0.42, 0.42, fill=RED, shape=MSO_SHAPE.OVAL)
        textbox(slide, cx, cy + 0.09, 0.42, 0.26, [num], 11.5, WHITE, bold=True,
                align=PP_ALIGN.CENTER, font=FONT_EN)
        cy += num_h
    cy += max(0.0, inner - (num_h + th + tgap + bh)) / 2
    if title:
        textbox(slide, cx, cy, cw, th + 0.04, [title], tp,
                RED if accent else INK, bold=True)
        cy += th + tgap
    if bodies:
        textbox(slide, cx, cy, cw, bh + 0.04, bodies, bp, INK,
                lsp=BODY_LSP, gap=BODY_GAP)


# ------------------------------------------------------------- slide kinds
def s_cover(prs, d):
    s = new_slide(prs)
    brand(s)
    for cx, cy, cd, col in ((10.35, 1.20, 0.80, RED),
                            (11.60, 2.20, 1.90, GREY_BG),
                            (9.85, 3.95, 0.36, RED),
                            (9.30, 1.95, 0.70, GREY_BG)):
        rect(s, cx, cy, cd, cd, fill=col, shape=MSO_SHAPE.OVAL)
    rect(s, ML, 1.28, 0.95, 0.95, fill=CARD, line=LINE, shape=MSO_SHAPE.OVAL)
    rect(s, ML + 0.34, 1.62, 0.28, 0.28, fill=RED)
    kicker(s, ML, 2.62, 8.2, d["kicker"], pt=12)
    textbox(s, ML, 3.00, 8.6, 1.96, d["title"], 48, INK, bold=True, lsp=0.94)
    textbox(s, ML, 5.06, 8.0, 0.90, d["sub"], 14.5, MUTED, lsp=1.06)
    rect(s, 0, 6.30, W, 1.20, fill=RED)
    textbox(s, ML, 6.60, 9.6, 0.66, d["meta"], 11.5, WHITE, lsp=1.06)
    return s


def s_agenda(prs, d):
    s = new_slide(prs)
    brand(s)
    top, bottom = 1.05, 6.85
    mid = (top + bottom) / 2
    rect(s, ML, mid - 1.70, 4.05, 3.40, fill=RED)
    kicker(s, ML + 0.42, mid - 1.22, 3.2, "CONTENT", WHITE, pt=11)
    textbox(s, ML + 0.42, mid - 0.86, 3.2, 0.92, ["目录"], 40, WHITE, bold=True)
    textbox(s, ML + 0.42, mid + 0.34, 3.2, 0.30,
            [f"{len(d['items'])} SECTIONS"], 12, WHITE, bold=True, font=FONT_EN)
    x, w = ML + 4.45, CW - 4.45
    n = len(d["items"])
    gap = 0.16
    hh = (bottom - top - gap * (n - 1)) / n
    for i, t in enumerate(d["items"]):
        y = top + i * (hh + gap)
        rect(s, x, y, w, hh, fill=GREY_BG)
        rect(s, x, y, 0.66, hh, fill=RED)
        textbox(s, x, y, 0.66, hh, [f"{i + 1:02d}"], 12, WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_EN)
        autotext(s, x + 0.90, y, w - 1.20, hh, [t], 15, INK,
                 anchor=MSO_ANCHOR.MIDDLE)
    return s


def s_section(prs, d):
    s = new_slide(prs)
    brand(s)
    rect(s, ML, 2.06, 0.98, 0.98, fill=None, line=LINE, lw=1.4,
         shape=MSO_SHAPE.OVAL)
    textbox(s, ML, 2.38, 0.98, 0.40, [d["no"]], 21, RED, bold=True,
            align=PP_ALIGN.CENTER, font=FONT_EN)
    kicker(s, ML, 3.26, 4.4, d["kicker"], pt=11)
    textbox(s, ML, 3.60, 5.4, 1.56, d["title"], 38, INK, bold=True, lsp=0.98)
    autotext(s, ML, 5.30, 5.0, 0.86, [d["lead"]], 14, MUTED, lsp=1.06)
    rect(s, 6.55, 2.20, W - ML - 6.55, 3.10, fill=RED)
    kicker(s, 6.95, 2.62, 4.0, d["key_label"], WHITE, pt=11)
    autotext(s, 6.95, 3.02, W - ML - 6.95 - 0.40, 1.90, [d["key"]], 20, WHITE,
             bold=True, floor=12, lsp=1.04)
    return s


def s_closing(prs, d):
    s = new_slide(prs, bg=DARK)
    brand(s, on_dark=True)
    rect(s, ML, 1.90, 0.95, 0.95, fill=RGBColor(0x22, 0x26, 0x2D),
         shape=MSO_SHAPE.OVAL)
    rect(s, ML + 0.34, 2.24, 0.28, 0.28, fill=RED)
    kicker(s, ML, 3.20, 4.0, "THANK YOU", pt=12)
    textbox(s, ML, 3.58, 9.0, 1.80, d["title"], 44, WHITE, bold=True, lsp=0.96)
    textbox(s, ML, 5.60, 8.6, 0.90, d["lead"], 14,
            RGBColor(0x9A, 0xA3, 0xAF), lsp=1.10)
    return s


def s_media(prs, d):
    """Big stat column + figure on top, three takeaway cards underneath."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.06
    cards = d.get("cards", [])
    cwid = (CW - 2 * 0.26) / 3
    ch = row_h(cards, cwid, 13, 11, 0.22, 2.45, floor_h=1.5) if cards else 0.0
    fig_h = (BOT - y - ch - 0.78) if cards else (BOT - y - 0.40)
    sw = 2.60
    kicker(s, ML, y, sw, d["kicker"], pt=10.5)
    textbox(s, ML, y + 0.34, sw, 0.78, [d["stat"]], 40, RED, bold=True,
            font=FONT_EN)
    autotext(s, ML, y + 1.20, sw, 0.72, d["stat_sub"], 11.5, MUTED, lsp=1.06)
    rule(s, ML, y + 2.04, 0.90)
    fx = ML + sw + 0.42
    fw = W - ML - fx
    _, _, _, dh = image(s, IMG / d["img"], fx, y, fw, fig_h)
    autotext(s, fx, y + (fig_h + dh) / 2 + 0.10, fw, 0.32, [d["caption"]], 10,
             MUTED, floor=8)
    if cards:
        cy = BOT - ch
        for i, c in enumerate(cards):
            card_text(s, ML + i * (cwid + 0.26), cy, cwid, ch, c["t"], c["b"],
                      accent=c.get("accent", False), flat=True, tpt=13,
                      bpt=11, pad=0.22)
    return s


def s_figure_side(prs, d):
    """Figure on the left, a stack of note cards / bands on the right."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.10
    fw = d.get("fig_w", 8.05)
    fh = BOT - y - 0.52
    _, dy, _, dh = image(s, IMG / d["img"], ML, y, fw, fh)
    autotext(s, ML, min(dy + dh + 0.12, BOT - 0.40), fw, 0.42, [d["caption"]],
             10, MUTED, floor=8)
    x = ML + fw + 0.42
    w = W - ML - x
    kicker(s, x, y, w, d["kicker"], pt=10.5)
    cy = y + 0.38
    blocks = d["blocks"]
    avail = BOT - cy - 0.16 * (len(blocks) - 1)
    weights = [b.get("wt", 1.0) for b in blocks]
    for b, wt in zip(blocks, weights):
        bh = avail * wt / sum(weights)
        if b["kind"] == "band":
            band(s, x, cy, w, bh, b["text"], base=12)
        else:
            card_text(s, x, cy, w, bh, b.get("t", ""), b["b"],
                      accent=b.get("accent", False), flat=True, tpt=12,
                      bpt=10.5, pad=0.20, floor=8)
        cy += bh + 0.16
    return s


def s_process(prs, d):
    """Four numbered steps in a row, optional closing band."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.10
    kicker(s, ML, y, CW, d["kicker"], pt=10.5)
    textbox(s, ML, y + 0.32, CW, 0.72, [d["title"]], 30, INK, bold=True)
    top = y + 1.24
    bh = 0.92 if d.get("band") else 0.0
    avail = BOT - top - (bh + 0.30 if bh else 0)
    n = len(d["steps"])
    gap = 0.24
    cwid = (CW - gap * (n - 1)) / n
    ch = row_h(d["steps"], cwid, 15, 11.5, 0.24, avail, num=True, floor_h=2.2)
    top += (avail - ch) / 2
    for i, st in enumerate(d["steps"]):
        card_text(s, ML + i * (cwid + gap), top, cwid, ch, st["t"], st["b"],
                  accent=st.get("accent", False), tpt=15, bpt=11.5,
                  num=f"{i + 1:02d}")
    if bh:
        band(s, ML, BOT - bh, CW, bh, d["band"], base=14)
    return s


def s_scheme(prs, d):
    """Scheme page: four mechanism steps on top, three result cards below."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.06
    kicker(s, ML, y, CW, d["kicker"], pt=10)
    textbox(s, ML, y + 0.28, CW, 0.62, [d["title"]], 26, INK, bold=True)
    top = y + 1.06
    n, gap = 4, 0.22
    cwid = (CW - gap * (n - 1)) / n
    ph = row_h(d["steps"], cwid, 14, 10.5, 0.20, 2.60, num=True, floor_h=2.0)
    for i, st in enumerate(d["steps"]):
        card_text(s, ML + i * (cwid + gap), top, cwid, ph, st["t"], st["b"],
                  accent=st.get("accent", False), tpt=14, bpt=10.5, pad=0.20,
                  num=f"{i + 1:02d}")
    cy = top + ph + 0.30
    ch = BOT - cy
    cwid = (CW - 2 * 0.24) / 3
    for i, c in enumerate(d["cards"]):
        card_text(s, ML + i * (cwid + 0.24), cy, cwid, ch, c["t"], c["b"],
                  accent=c.get("accent", False), flat=True, tpt=12.5, bpt=10.5,
                  pad=0.20)
    return s


def s_matrix(prs, d):
    """2 x 4 card matrix, optional closing band."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.10
    kicker(s, ML, y, CW, d["kicker"], pt=10.5)
    textbox(s, ML, y + 0.32, CW, 0.72, [d["title"]], 30, INK, bold=True)
    top = y + 1.28
    bh = 0.94 if d.get("band") else 0.0
    total = BOT - top - (bh + 0.26 if bh else 0)
    gx, gy = 0.24, 0.24
    cwid = (CW - 3 * gx) / 4
    rows = math.ceil(len(d["cells"]) / 4)
    tpt, bpt = (15, 12) if rows == 1 else (13.5, 10.5)
    ch = row_h(d["cells"], cwid, tpt, bpt, 0.22,
               (total - gy * (rows - 1)) / rows, floor_h=1.8)
    top += (total - (rows * ch + gy * (rows - 1))) / 2
    for i, c in enumerate(d["cells"]):
        r, col = divmod(i, 4)
        card_text(s, ML + col * (cwid + gx), top + r * (ch + gy), cwid, ch,
                  c["t"], c["b"], accent=c.get("accent", False), tpt=tpt,
                  bpt=bpt, pad=0.22)
    if bh:
        band(s, ML, BOT - bh, CW, bh, d["band"], base=13)
    return s


def s_triple(prs, d):
    """Three tall cards, optional closing band."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.10
    kicker(s, ML, y, CW, d["kicker"], pt=10.5)
    textbox(s, ML, y + 0.32, CW, 0.72, [d["title"]], 30, INK, bold=True)
    top = y + 1.32
    bh = 0.94 if d.get("band") else 0.0
    avail = BOT - top - (bh + 0.26 if bh else 0)
    gap = 0.26
    cwid = (CW - 2 * gap) / 3
    ch = row_h(d["cards"], cwid, 15, 11.5, 0.24, avail,
               num=any(c.get("num") for c in d["cards"]), floor_h=2.4)
    top += (avail - ch) / 2
    for i, c in enumerate(d["cards"]):
        card_text(s, ML + i * (cwid + gap), top, cwid, ch, c["t"], c["b"],
                  accent=c.get("accent", False), tpt=15, bpt=11.5,
                  num=c.get("num"))
    if bh:
        band(s, ML, BOT - bh, CW, bh, d["band"], base=13)
    return s


def s_bars(prs, d):
    """Per-core bandwidth bar chart drawn as shapes."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.36
    kicker(s, ML, y, 4.6, d["kicker"], pt=10.5)
    textbox(s, ML, y + 0.36, 5.0, 1.40, d["title"], 29, INK, bold=True, lsp=0.98)
    autotext(s, ML, y + 1.94, 4.9, 1.00, [d["lead"]], 13, MUTED, lsp=1.06)
    textbox(s, ML, y + 3.10, 4.0, 0.74, [d["stat"]], 38, RED, bold=True,
            font=FONT_EN)
    textbox(s, ML, y + 3.84, 4.6, 0.30, [d["stat_sub"]], 11, MUTED)

    bx = ML + 5.40
    bw = W - ML - bx
    card(s, bx, TOP + 0.20, bw, BOT - TOP - 0.20)
    ix, iy = bx + 0.34, TOP + 0.52
    iw, ih = bw - 0.68, BOT - TOP - 1.30
    rows = d["rows"]
    rh = ih / len(rows)
    track_x = ix + 0.62
    track_w = iw - 0.62 - 0.86
    for i, (name, val, frac) in enumerate(rows):
        cy = iy + i * rh
        bar_h = min(0.26, rh * 0.62)
        textbox(s, ix, cy + (rh - 0.22) / 2, 0.56, 0.22, [name], 11, MUTED,
                bold=True, font=FONT_EN)
        rect(s, track_x, cy + (rh - bar_h) / 2, track_w, bar_h, fill=GREY_BG)
        rect(s, track_x, cy + (rh - bar_h) / 2, track_w * frac, bar_h,
             fill=RED if frac > 0.8 else RED_DARK)
        textbox(s, ix + iw - 0.80, cy + (rh - 0.22) / 2, 0.80, 0.22, [val], 11,
                INK, bold=True, align=PP_ALIGN.RIGHT, font=FONT_EN)
    autotext(s, ix, iy + ih + 0.18, iw, 0.46, [d["caption"]], 10, MUTED,
             floor=8)
    return s


def s_compare(prs, d):
    """Expected vs measured, plus three stat tiles."""
    s = new_slide(prs)
    chrome(s, d["chrome"])
    brand(s)
    y = TOP + 0.06
    kicker(s, ML, y, CW, d["kicker"], pt=10.5)
    textbox(s, ML, y + 0.30, CW, 0.86, [d["title"]], 26, INK, bold=True,
            lsp=1.22)
    autotext(s, ML, y + 1.24, CW, 0.34, [d["lead"]], 13.5, MUTED)
    top = y + 1.70
    sh = 1.16
    ch = BOT - top - sh - 0.26
    cwid = (CW - 0.30) / 2
    for i, col in enumerate(d["cols"]):
        x = ML + i * (cwid + 0.30)
        card(s, x, top, cwid, ch, accent=col.get("accent", False))
        textbox(s, x + 0.28, top + 0.24, cwid - 0.56, 0.32, [col["t"]], 14,
                RED if col.get("accent") else MUTED, bold=True)
        autotext(s, x + 0.28, top + 0.72, cwid - 0.56, ch - 1.00, col["b"],
                 12, INK, floor=9, lsp=BODY_LSP, gap=10)
    ty = BOT - sh
    twid = (CW - 2 * 0.26) / 3
    for i, (big, sub) in enumerate(d["stats"]):
        x = ML + i * (twid + 0.26)
        card(s, x, ty, twid, sh, flat=True)
        textbox(s, x, ty + 0.16, twid, 0.56, [big], 28, RED, bold=True,
                align=PP_ALIGN.CENTER, font=FONT_EN)
        autotext(s, x + 0.22, ty + 0.74, twid - 0.44, 0.32, [sub], 10.5, MUTED,
                 floor=8, align=PP_ALIGN.CENTER)
    return s


# ------------------------------------------------------------------ live numbers
def _j2cov(j: float) -> float:
    j = min(max(float(j), 1e-12), 1.0)
    return math.sqrt((1.0 - j) / j)


class Live:
    """Numbers the deck is allowed to quote: only what deck_ring2_data.json ran."""

    def __init__(self, d: dict):
        self.d = d
        self.oc = int(d["meta"]["core_outstanding"])
        self.oc16 = int(d["meta"].get("s16_overcommit") or 20)
        self.kw = int(d["meta"]["k_write"])
        self.kr = int(d["meta"]["k_read"])
        self.rstar = float(d["ideal"]["r_fair"])
        self.rmax = float(d["ideal"]["r_max"])
        met = json.loads((ROOT / "results" / "metric_ring2_cc.json").read_text())
        self.kappa = float(met["kappa"])

    def W(self, name: str) -> dict:
        return self.d["write"][name]

    def R(self, name: str) -> dict:
        return self.d["read"][name]

    def P(self, name: str) -> dict:
        return self.d["read_payload"][name]

    def thr(self, name: str, nd: int = 4) -> str:
        return f"{self.W(name)['throughput']:.{nd}f}"

    def covf(self, name: str) -> float:
        return _j2cov(self.W(name)["jain_bin"]["jain_bin_mean"])

    def cov(self, name: str, nd: int = 4) -> str:
        return f"{self.covf(name):.{nd}f}"

    def mm(self, name: str, nd: int = 4) -> str:
        return f"{self.W(name)['max_min']:.{nd}f}"

    def pct_r(self, name: str, nd: int = 2) -> str:
        return f"{100.0 * self.W(name)['throughput'] / self.rstar:.{nd}f}"

    def dthr(self, name: str, vs: str = "S0", nd: int = 1) -> str:
        a, b = self.W(name)["throughput"], self.W(vs)["throughput"]
        return f"{100.0 * (a - b) / b:+.{nd}f}%"

    def dcov(self, name: str, vs: str = "S0", nd: int = 3) -> str:
        return f"{self.covf(name) - self.covf(vs):+.{nd}f}"

    def phif(self, name: str) -> float:
        return (self.W(name)["throughput"] - self.kappa * self.covf(name)) / self.rstar

    def phi(self, name: str, nd: int = 3) -> str:
        return f"{self.phif(name):.{nd}f}"

    def triple(self, name: str) -> str:
        return f"{self.thr(name, 3)} / {self.cov(name)} / {self.mm(name)}"

    def bw(self, name: str, core: int, nd: int = 5) -> str:
        return f"{float(self.W(name)['bw_by_core'][str(core)]):.{nd}f}"

    def dbw(self, name: str, core: int, vs: str = "S0", nd: int = 1) -> str:
        a = float(self.W(name)["bw_by_core"][str(core)])
        b = float(self.W(vs)["bw_by_core"][str(core)])
        return f"{100.0 * (a - b) / b:+.{nd}f}%"

    def finish(self, name: str, core: int) -> int:
        return int(self.W(name)["finish_by_core"][str(core)])

    def finish_span(self, name: str) -> tuple[int, int]:
        fin = self.W(name)["finish_by_core"]
        vs = [int(v) for v in fin.values()]
        return min(vs), max(vs)

    def fratf(self, name: str) -> float:
        lo, hi = self.finish_span(name)
        return hi / lo

    def frat(self, name: str, nd: int = 3) -> str:
        return f"{self.fratf(name):.{nd}f}"

    def fspan(self, name: str) -> str:
        lo, hi = self.finish_span(name)
        return f"{lo:,}–{hi:,}"

    def cg(self, name: str) -> dict:
        return self.W(name)["ceiling_gap"]

    def util(self, name: str, nd: int = 2) -> str:
        return f"{100.0 * float(self.cg(name)['util']):.{nd}f}%"

    def floor_cov(self, name: str, nd: int = 4) -> str:
        return f"{_j2cov(self.W(name)['regular']['jain_regular']):.{nd}f}"

    def p05_cov(self, name: str, nd: int = 3) -> str:
        return f"{_j2cov(self.W(name)['jain_bin']['jain_bin_p05']):.{nd}f}"

    def min_cov(self, name: str, nd: int = 3) -> str:
        return f"{_j2cov(self.W(name)['jain_bin']['jain_bin_min']):.{nd}f}"

    def nbins(self, name: str) -> int:
        return int(self.W(name)["jain_bin"]["n_bins"])

    def rcov(self, name: str, nd: int = 4) -> str:
        return f"{_j2cov(self.R(name)['jain_bin']['jain_bin_mean']):.{nd}f}"

    def rmm(self, name: str, nd: int = 4) -> str:
        return f"{self.R(name)['max_min']:.{nd}f}"

    def rthr(self, name: str, nd: int = 4) -> str:
        return f"{self.R(name)['throughput']:.{nd}f}"

    def rpct(self, name: str, nd: int = 2) -> str:
        return f"{100.0 * self.R(name)['throughput'] / float(self.d['ideal']['read_r_fair']):.{nd}f}"

    def rdthr(self, name: str, vs: str = "S0", nd: int = 1) -> str:
        a, b = self.R(name)["throughput"], self.R(vs)["throughput"]
        return f"{100.0 * (a - b) / b:+.{nd}f}%"

    def pcov(self, name: str, nd: int = 4) -> str:
        return f"{_j2cov(self.P(name)['jain_bin']['jain_bin_mean']):.{nd}f}"

    def pmm(self, name: str, nd: int = 4) -> str:
        return f"{self.P(name)['max_min']:.{nd}f}"

    def pthr(self, name: str, nd: int = 4) -> str:
        return f"{self.P(name)['throughput']:.{nd}f}"

    def pfrat(self, name: str, nd: int = 3) -> str:
        fin = [int(v) for v in self.P(name)["finish_by_core"].values()]
        return f"{max(fin) / min(fin):.{nd}f}"

    def pfspan(self, name: str) -> str:
        fin = [int(v) for v in self.P(name)["finish_by_core"].values()]
        return f"{min(fin):,}–{max(fin):,}"

    def down_fail(self, name: str) -> int:
        w = self.W(name)
        return int(w.get("n_down_fail") or (
            int(w.get("n_deflections") or 0)
            + int(w.get("n_leave_occ_gt1") or 0)))

    def leave_occ(self, name: str) -> int:
        return int(self.W(name).get("n_leave_occ_gt1") or 0)

    def same_op(self, a: str, b: str) -> bool:
        wa, wb = self.W(a), self.W(b)
        return (abs(float(wa["throughput"]) - float(wb["throughput"])) < 5e-4
                and int(wa.get("makespan") or 0) == int(wb.get("makespan") or 0))


def slides(n: Live) -> list:
    s0_lo, s0_hi = n.finish_span("S0")
    s1_lo, s1_hi = n.finish_span("S1")
    s0_bw = {int(c): float(v) for c, v in n.W("S0")["bw_by_core"].items()}
    s0_slow = min(s0_bw, key=s0_bw.get)
    s0_fast = max(s0_bw, key=s0_bw.get)
    s0_rows = [(f"C{c}", f"{s0_bw[c]:.3f}", s0_bw[c] / 0.75)
               for c in sorted(s0_bw)]
    s0g = n.cg("S0")
    s0reg = n.W("S0")["regular"]
    s1_fail = int(n.W("S1")["n_board_fail"])
    s1_defl = n.down_fail("S1")
    s1_occ = n.leave_occ("S1")
    s1_ring_defl = int(n.W("S1").get("n_deflections") or 0)
    s1d_eq = n.same_op("S1D", "S0")
    s1u_eq = n.same_op("S1U", "S1")
    gap_r = 100.0 - float(n.pct_r("S0"))
    return [
    ("cover", dict(
        kicker="NOC ARCHITECTURE REVIEW · 20-NODE BUFFERLESS RING",
        title=["满带宽下的", "多核带宽不均"],
        sub=["研究问题：无缓存环接近 R* 时，几何造成的六快四慢",
             "能否在不牺牲总带宽的前提下被控制器抹平？"],
        meta=["汇报对象：芯片架构师",
              f"口径：问题 → 度量 → 机制 → 上界 → 决策 · "
              f"uniform tiled 写 · K = {n.kw} · "
              f"每核 outstanding = {n.oc} · 主指标 100 拍窗 CoV"])),

    ("agenda", dict(items=[
        "问题与口径：拓扑、争用窗、100 拍 CoV、outstanding = 128",
        "S0：现象 → 逐拍账本 → 结构根因（几何 × 在环优先）",
        "S1：机制、实测、信号消融 —— 为何只能换一端",
        "设计空间：信号 × 控制点，空着的四类补齐并同口径实测",
        "理论上界 R(CoV) 与单一标量 φ = (R − κ·CoV)/R*",
        "可实现点：S16 / S22 / S29 的机制、硬件与实测",
        "结论：一项架构决策（HA 授权能不能改）"])),

    ("section", dict(
        no="01", kicker="SECTION ONE", title=["问题背景"],
        lead="拓扑、节点角色、路由规则，以及决定结论的那几项关键 setup。",
        key_label="KEY MESSAGE",
        key="10 个 AI core 对 8 个 memory HA 做 CHI WriteNoSnp。"
            "先固定度量（争用窗 + 100 拍 CoV），再问总带宽和瞬时份额各卡在哪。")),

    ("figside", dict(
        chrome="拓扑与节点", img="04-topo.png", fig_w=5.90,
        caption="蓝 = AI core（10）· 橙 = memory HA（8）· 灰 = 非终端 N9 / N19；"
                "边上数字为 hop 时延（拍）。",
        kicker="TOPOLOGY & ROUTING",
        blocks=[
            dict(kind="card", t="20 节点单 plane 双向闭合 full ring",
                 b=["**节点角色**：10 个 AI core（C0–C18 偶数）发起写；"
                    "8 个 memory HA（M1–M17 奇数）为 completer；N9 / N19 只转发，不收发。",
                    "**路由：时延最短**。按链路时延之和选 CW / CCW，时延平局比跳数，"
                    "再平局取 CW；core→mem 与 mem→core 同一规则。"], wt=1.25),
            dict(kind="card", accent=True, t="几何后果：本次问题的起点",
                 b=["因为 N9 / N19 不是 mem，C0 / C8 / C10 / C18 的「邻接 mem 数」只有 1，"
                    "其余六核为 2。这一条几何差异，直接决定了后面看到的六快四慢。"], wt=1.0),
            dict(kind="band",
                 text="20 节点 / 1 plane · 10 core + 8 HA + 2 非终端 · "
                      "等速率理论上限 R* = 5.7143 flit/cycle", wt=0.62)])),

    ("process", dict(
        chrome="协议：CHI WriteNoSnp", kicker="PROTOCOL",
        title="一笔写 = 四拍握手，被计量的只有第三拍",
        steps=[
            dict(t="REQ", b=["core → HA，1 flit，走时延最短方向，占 REQ VC。"]),
            dict(t="DBIDResp",
                 b=["HA → core，1 flit，反方向，占 RSP VC。HA think time = 0。"]),
            dict(t="WriteData × 2", accent=True,
                 b=["core → HA，2 flit，占 DAT VC。",
                    "**per-core 写带宽 = 争用窗内成功上环的 WriteData flit / cycle。**"]),
            dict(t="Comp",
                 b=["HA → core，1 flit，占 RSP VC，释放该笔事务的 tracker 表项。"])],
        band="度量：闭环批量每核 K 笔，争用窗 [0, 第一核耗尽]；"
             "主指标 = 100 拍窗内十核写带宽的 CoV（标准差 / 均值）再对窗取平均；"
             "辅指标 = 全程总带宽 / R* 与各核完成时间。全篇同一口径。")),

    ("matrix", dict(
        chrome="关键 setup", kicker="FABRIC CONFIG", title="八项决定结论的配置",
        cells=[
            dict(t="上环端口",
                 b=["每方向一组 × REQ / RSP / DAT，每口 1 flit/cycle → 每节点 6 个上环口。"]),
            dict(t="下环端口",
                 b=["每 (node, VC) 一个「两写一读」buffer：两方向可同拍各写 1 flit，"
                    "PE 每拍只读 1 flit。"]),
            dict(t="上环队列",
                 b=["每 (node, VC) 12 深共享 FIFO + 每向 8 深 inject Q，各接一个上环口。"]),
            dict(t="下环队列",
                 b=["深度 12。整轮里满的比例 0%，从来不是瓶颈 —— "
                    "所以带宽不是被下环侧卡住的。"]),
            dict(t="在环绝对优先", accent=True,
                 b=["bufferless：transit flit 永不被本地注入打断，注入只能挤空隙。"]),
            dict(t="I-tag / E-tag",
                 b=["t_inj = 16、hold = 8 定向上游预约单槽；t_xfer = 1，首次下环失败即打标，"
                    "再到达时最高优先级下环。"]),
            dict(t=f"每核 outstanding = {n.oc}", accent=True,
                 b=["全文统一设定，不是扫描轴。最坏空载写 RTT 89 拍 × 等速率 2/7 笔/拍 "
                    f"≈ 26 笔在飞，所以 {n.oc} 不绑定环；"
                    f"HA tracker 512 远大于 10 × {n.oc} 的供给，也不绑定。"]),
            dict(t="专用流控总线",
                 b=["广播，不占 NoC hop；**任何使用固定 30 拍延迟**（物理约束，"
                    "不是可调项），控制窗 64 拍。"])])),

    ("section", dict(
        no="02", kicker="SECTION TWO", title=["S0 的具体问题"],
        lead=f"总带宽已经打到理论上限的 {n.pct_r('S0')}%，"
             f"但十个核的长期速率差 {n.mm('S0')} 倍。",
        key_label="KEY MESSAGE",
        key="带宽这条线基本关死了；真正的缺陷在公平性一侧，"
            "而且是几何决定的固定分组，不是随机抖动。")),

    ("media", dict(
        chrome="S0：总带宽已打满", kicker="BANDWIDTH IS SATURATED",
        stat=f"{n.pct_r('S0')}%",
        stat_sub=[f"S0 = {n.thr('S0')} flit/cycle",
                  f"占等速率上限 R* = {n.rstar:.4f}"],
        img="07-totalbw.png",
        caption="左 = 每个 VC 的 40 条有向链路按占用率排序，天花板是「每拍 1 flit」；"
                "右 = 绑定链路的拍数预算。两幅都不把分箱带宽和理论上限放在同一根轴上。",
        cards=[
            dict(t="判据：链路已经快满了",
                 b=[f"最忙 hop（{s0g['hop']}）占用 **{n.util('S0')}**。"
                    "它们就是 R* 的绑定项 —— "
                    "带宽是被链路卡住的，不是被队列或端口卡住的。"]),
            dict(t=f"差的 {gap_r:.2f}% 是什么", accent=True,
                 b=[f"makespan {s0g['makespan']:,} = {s0g['floor']:,} 载有效载荷 + "
                    f"{s0g['surcharge']:,} 绕环重发 + **{s0g['idle']:,} 空转**，"
                    "三项逐拍相加相等。"
                    f"core_outstanding 固定为 {n.oc}、覆盖最坏 RTT；在这个前提下，"
                    "环没有被在飞上限卡住。"]),
            dict(t="官方 S0 已是 I-tag 的带宽最高点",
                 b=["I-tag 工作点 t_inj = 16、hold = 8 来自官方网格的带宽最优，"
                    f"总带宽 {n.pct_r('S0')}% R*。"
                    "再往公平一侧调（更短 t_inj / 更短 hold）会掉带宽；"
                    "加深下环队列也不能抬高这条绑定 hop。"])])),

    ("bars", dict(
        chrome="S0：各核写带宽", kicker="PER-CORE BANDWIDTH",
        title=["六快四慢，", "是固定分组不是抖动"],
        lead=f"最低 C{s0_slow} = {s0_bw[s0_slow]:.5f}，最高 C{s0_fast} = {s0_bw[s0_fast]:.5f}。"
             "慢的那四个正是「邻接 mem 数 = 1」的 C0 / C8 / C10 / C18。",
        stat=n.mm("S0"), stat_sub=f"整窗 MAX / MIN（S0，K = {n.kw}）",
        rows=s0_rows,
        caption="条长按 0.75 flit/cycle 归一。8 个 HA 收到的 WriteData 完全相等"
                "（50,000 / HA），所以这不是访存不均。")),

    ("media", dict(
        chrome="S0：100 拍瞬时不均衡度", kicker="INSTANTANEOUS IMBALANCE",
        stat=n.cov("S0"),
        stat_sub=["100 拍窗十核带宽的 CoV（标准差 / 均值）", "理想控制器 ≈ 0（整数粒度）"],
        img="09-s0-instbal.png",
        caption="左 = 主指标随时间；右 = 同一份数据换观察窗宽度，实测始终高于理想控制器。",
        cards=[
            dict(t="不是个别坏箱",
                 b=[f"{n.nbins('S0')} 个箱里最差 5% 的箱 CoV ≥ {n.p05_cov('S0')}，"
                    f"最差箱 {n.min_cov('S0')}。右图把观察窗一路放宽："
                    f"抹平抖动后仍收敛到 {n.floor_cov('S0')} 的地板，到不了理想控制器。"]),
            dict(t=f"抖动抹平后的下限仍有 {n.floor_cov('S0')}", accent=True,
                 b=["给每核每箱都填上它自己的长期均值（抖动全抹掉、只留速率差），"
                    "CoV 就降到底了 —— 这是任何「只整时机、不搬份额」机制的地板。"]),
            dict(t="所以要降低不均衡度，必须真的搬速率",
                 b=[f"核间长期速率 {s0reg['rate_min']:.2f} vs {s0reg['rate_max']:.2f} "
                    f"flit/箱（比 {s0reg['rate_ratio']:.4f}）。"
                    f"核内抖动占方差 {100*s0reg['within_share']:.1f}%，"
                    f"但封顶的是另外那 {100*(1-s0reg['within_share']):.1f}% 的核间差 —— "
                    f"100 拍窗内每核有 {s0reg['flits_per_core_per_bin']:.1f} 个 flit，"
                    "整数粒度不构成限制。"])])),

    ("triple", dict(
        chrome="根因", kicker="ROOT CAUSE", title="几何 × 机制 = 固定的六快四慢",
        cards=[
            dict(num="01", t="几何：邻接 mem = 1",
                 b=["N9 / N19 是非终端，于是 C0 / C8 / C10 / C18 紧邻的 memory 只有一个。"
                    "它们的写必须挤到环上最忙的那几条 hop 上，而其余六核可以在两侧分流。",
                    "[[带宽与邻接 mem 数的相关系数 0.997。改路由也变不出第二个邻居。]]"]),
            dict(num="02", t="机制：在环绝对优先",
                 b=["无缓存环上 transit flit 永不被本地注入打断，注入只能挤空隙。"
                    "决定上环延迟的是空隙的**分布**而不是总量 —— "
                    "坐在热段起点的节点看到的空隙最少、最不规则。",
                    f"[[官方 run 最忙 hop 占用 {n.util('S0')}；短探测把绑定 hop 的"
                    "空槽逐拍分类，「端口空闲」一列合计近零 —— "
                    "没有一个空槽是仲裁失误。]]"]),
            dict(num="03", t="症状：差距全在「等着上环」这一段",
                 b=["8 个 HA 收到的 WriteData 完全相等，每核每方向**成功**上环数也被路由"
                    "钉死在 30,000 笔、逐位相同。服务侧没有任何不均。",
                    "[[所以差的不是「谁被服务得多」，而是「谁等得久」—— "
                    "决策点在注入排队，不在存储侧。]]"])])),

    ("section", dict(
        no="03", kicker="SECTION THREE", title=["S1 方案与实测"],
        lead="拥塞检测、拥塞传递、拥塞反馈、流量控制，以及两个工作点的实测。",
        key_label="KEY MESSAGE",
        key=f"S1 默认档：CoV {n.cov('S0')} → {n.cov('S1')}（{n.dcov('S1')}），"
            f"总写带宽 {n.dthr('S1')}；"
            "S1T 把带宽收回后，CoV 回到与 S0 同一水平。"
            + ("计入 leave FIFO 占用 > 1 之后，S1D 不再等于 S0："
               "下环一路能驱动等级，但几乎只掉带宽、不搬公平。"
               if not s1d_eq else
               "下环失败含 FIFO 占用 > 1 后次数变大，核侧 64 拍窗仍到不了等级 1。"))),

    ("process", dict(
        chrome="S1 机制：检测 / 传递 / 反馈 / 控制",
        kicker="SENDER-DRIVEN · RATE-BASED · BUS-TRIGGERED",
        title="S1：拥塞等级广播 + 节点级 AIMD",
        steps=[
            dict(t="拥塞检测",
                 b=["每节点每 64 拍窗统计**上环失败**与**下环失败**"
                    "（绕环偏转 + leave FIFO 写入后占用 > 1），"
                    "分 total（任何原因）与 net（仅在环占用造成）两路。",
                    "等级 = min(7, 计数 ÷ 8)，量化成 3 bit：0–7→0，8–15→1，…，≥56→7。"]),
            dict(t="拥塞传递",
                 b=["每方向 3 bit 的**专用广播总线**，从不占用 NoC hop，"
                    "因此不会自我加剧拥塞。",
                    "本设计里任何总线使用固定 **30 拍**延迟（物理约束）；30 < 窗口 64，"
                    "所以反馈仍能赶在下一次 AIMD 之前送到。"]),
            dict(t="拥塞反馈",
                 b=["每节点维护一张**受控节点表** —— 自己的 flit 会经过的那些 "
                    "path nodes（20 项 × 6 bit）。",
                    "反馈值取这些节点等级的 **max**：路径上最堵的那一段说了算。"]),
            dict(t="流量控制", accent=True,
                 b=["最终等级 = level_of(自身 total 失败 − 8 × 收到的 max net 等级)。",
                    "对每窗**整型注入预算**做 AIMD：乘性减 α = 0.75 / 0.5 / 0.25，"
                    "加性增 β = +16 / +8 / +2，出口是令牌桶闸门。"])],
        band="第 ① 段的信号是本节点上环失败。无缓存环上这类失败多数来自他人 transit，"
             "因此等级同时反映「本核注入多」和「本核被路过」。")),

    ("media", dict(
        chrome="S1 效果：默认档与调参档", kicker="TWO OPERATING POINTS",
        stat=n.dthr("S1"), stat_sub=["S1 默认档相对 S0 的总写带宽",
                                     f"CoV {n.dcov('S1')}"],
        img="12-s1-effect.png",
        caption="左 = 每核写带宽对比；右 = 三个工作点在带宽—CoV 平面上的位置。"
                "下环失败含 leave FIFO 占用 > 1。",
        cards=[
            dict(t="S1 默认（band = spec）",
                 b=[f"100 拍窗 CoV {n.cov('S0')} → **{n.cov('S1')}**，整窗 max/min "
                    f"{n.mm('S0')} → {n.mm('S1')}；"
                    f"总写带宽 {n.thr('S0')} → **{n.thr('S1')}（{n.dthr('S1')}）**。"]),
            dict(t="S1T 每向预算（调参后）", accent=True,
                 b=["62 组 AIMD 参数网格选出的最优点（dir_split、cap 0.5、w 64、"
                    f"burst 1）：带宽 **{n.thr('S1T')}（{n.dthr('S1T')}）**，"
                    f"CoV **{n.cov('S1T')}**，max/min {n.mm('S1T')}。"]),
            dict(t="两档各取一端",
                 b=[f"从 S1 走到 S1T：带宽 {n.dthr('S1T', 'S1')}，CoV {n.dcov('S1T', 'S1')}。"
                    "带宽接近 S0 的配置，CoV 也回到 S0 附近。"])])),

    ("compare", dict(
        chrome="S1 实测：各核带宽如何变化", kicker="MEASURED PER-CORE SHIFT",
        title="快核被压下来，慢核窗内速率略升，总带宽仍掉",
        lead=f"max/min 从 {n.mm('S0')} 降到 {n.mm('S1')}：差距在收，总量也在收。",
        cols=[
            dict(t="机制在做什么", b=[
                "**检测**：本节点上环失败 / 下环失败（真偏转 + leave FIFO 占用 > 1），"
                "量化成 3 bit 拥塞等级。",
                "**控制**：按路径上最高等级对注入预算做 AIMD，出口是令牌桶。",
                "**执行**：没额度的拍不上环。"]),
            dict(t="实测各核变化", accent=True, b=[
                "**信号来源**：乘性减仍主要由本节点上环失败触发；"
                "下环一路现在把 leave FIFO 占用 > 1 也算进去。",
                f"**快核被压、慢核略升**：慢核 C8 {n.bw('S0', 8)} → {n.bw('S1', 8)}"
                f"（{n.dbw('S1', 8)}），快核 C14 {n.bw('S0', 14)} → {n.bw('S1', 14)}"
                f"（{n.dbw('S1', 14)}）。",
                "**令牌桶**：没额度就不上环；该槽沿途节点都可以使用。"])],
        stats=[(n.dthr("S1"), "S1 默认档相对 S0 的总写带宽"),
               ("62 组", "AIMD 参数扫描的工作点数"),
               ("21,220", "S1 硬件 FF‑eq（总线 + 表 + 乘法器 + 令牌桶）")])),

    ("media", dict(
        chrome="S1 信号细分：只计上环失败 / 只计下环失败 / 两者都计",
        kicker="S1 SIGNAL SPLIT",
        stat=f"{s1_defl:,} / {s1_fail:,}",
        stat_sub=["S1 全程下环失败 / 上环失败",
                  f"下环 = 真偏转 {s1_ring_defl:,} + FIFO 占用>1  {s1_occ:,}"],
        img="38-s1-signal.png",
        caption="四根柱 = S0 / S1D（只计下环）/ S1U（只计上环）/ S1（都计），K = 20000 uniform 写；"
                "下环失败 = 绕环偏转 + 两写一读 leave FIFO 占用 > 1。",
        cards=[
            dict(t="三档怎么实现",
                 b=["参数 signal = up / down / both：**up** 只计上环失败；"
                    "**down** 只计下环失败（真偏转，以及 leave FIFO 写入后深度 > 1）；"
                    "**both** = 现有 S1，取二者 max。FIFO 本身仍是两写一读、深度 12。"]),
            dict(t=("实测：两档仍各与一端重合" if (s1d_eq and s1u_eq)
                    else "实测：计入 FIFO 占用后下环一路不再沉默"),
                 accent=True,
                 b=[f"**S1U{' = S1' if s1u_eq else ''}**：总带宽 {n.thr('S1U')}、"
                    f"100 拍 CoV {n.cov('S1U')}、完成时间 {n.fspan('S1U')}。",
                    f"**S1D{' = S0' if s1d_eq else ''}**：{n.thr('S1D')} / "
                    f"{n.cov('S1D')} / {n.fspan('S1D')}。",
                    ("四组曲线两两完全一致：全局下环次数变大，但核侧 64 拍窗仍到不了等级 1。"
                     if (s1d_eq and s1u_eq)
                     else "S1D 不再等于 S0：leave FIFO 占用 > 1 已经能驱动 down 等级。")]),
            dict(t="新口径下的量级",
                 b=[f"S1 全程上环失败 **{s1_fail:,}** 次；下环失败 **{s1_defl:,}** 次"
                    f"（其中绕环偏转 {s1_ring_defl:,}，FIFO 占用>1  {s1_occ:,}）。",
                    "[[占用 > 1 只改计数与 S1 的 down 信号，不把 flit 再送回环，也不打 E-tag。]]"])])),

    ("media", dict(
        chrome="S1 信号细分：各核完成时间曲线",
        kicker="PER-CORE COMPLETION",
        stat=f"{n.frat('S1D')} → {n.frat('S1U')}",
        stat_sub=[("十核完成时间 最晚 / 最早：S1D（= S0）→ S1U（= S1）"
                   if (s1d_eq and s1u_eq)
                   else "十核完成时间 最晚 / 最早：S1D → S1U"),
                  "同一组四个慢核，差距收窄但仍在"],
        img="39-s1-signal-finish.png",
        caption="横轴 cycle，纵轴 = 该核已上环 WriteData / 本核配额；红 = 邻 mem = 1 的四核"
                "（C0 / C8 / C10 / C18），蓝 = 其余六核；右 = 十核最后一个 flit 上环的时刻。"
                "下环失败含 leave FIFO 占用 > 1。",
        cards=[
            dict(t=("S1D（= S0）：六快四慢" if s1d_eq else "S1D：只听下环失败"),
                 b=[f"十核完成时间 {n.fspan('S1D')}，比值 {n.frat('S1D')}。"
                    + ("快核完成的瞬间慢核斜率抬升 —— "
                       "慢核不是能力不足，是满载时抢不到上环机会。"
                       if s1d_eq else
                       "leave FIFO 占用 > 1 已经进入 down 信号，完成曲线不再等于 S0。")]),
            dict(t=("S1U（= S1）：整体推后、差距收窄" if s1u_eq
                    else "S1U：只听上环失败"),
                 accent=True,
                 b=[f"十核完成时间 {n.fspan('S1U')}，比值 {n.frat('S1U')}；"
                    f"最早完成时刻从 {s0_lo:,} 推到 {n.finish_span('S1U')[0]:,}。"]),
            dict(t="两个口径说的不是一件事",
                 b=[f"100 拍 CoV 看窗内瞬时份额，完成时间比看长期速率。"
                    f"S1 把 CoV 从 {n.cov('S0')} 降到 {n.cov('S1')}（{n.dcov('S1')}），"
                    f"完成时间比从 {n.frat('S0')} 到 {n.frat('S1')}，"
                    "四慢核依然是同一组核。全部 13 个方案的完成曲线见第六节末两页。"])])),

    ("section", dict(
        no="04", kicker="SECTION FOUR", title=["现有拥塞控制", "算法的分类"],
        lead="按第 21 页的两个正交维度分类：先按拥塞信号（行），再按控制点（列），"
             "给出各类优缺点，定位 S1，并补齐本研究原本没有代表方案的四个格子。",
        key_label="KEY MESSAGE",
        key="S1 落在「源端速率 × 显式等级」：降速由本节点失败触发。"
            "无缓存环上这类失败多数来自他人 transit。同一行的方案共用这一信号。")),

    ("matrix", dict(
        chrome="分类 · 拥塞信号（第 21 页的行）",
        kicker="TAXONOMY BY SIGNAL  =  PAGE 19 ROWS",
        title="听什么决定盲区（= 第 21 页的行）",
        cells=[
            dict(t="无信号 / 1 bit 需求", b=[
                "**代表**：S29。",
                "**听什么**：固定日历 + 每核 1 bit「有活」。",
                "**优**：无测量、无收敛。**缺**：不知道谁落后。"]),
            dict(t="本地观测", b=[
                "**代表**：S16 / S21 / S26 / I-tag。",
                "**听什么**：本节点的占用、失败或授权数。",
                "**优**：零新报文。**缺**：源端格子会把别人造成的失败算到自己头上。"]),
            dict(t="链路占用", b=[
                "**代表**：S27。",
                "**听什么**：hop 占用率，沿邻线上传。",
                "**优**：信号就在拥塞 hop。**缺**：无缓存环上「保护链路」= 空转。"]),
            dict(t="时延", b=[
                "**代表**：S17 / S19。",
                "**听什么**：带内 RTT（Comp 回程）。",
                "**优**：不需要专用总线。**缺**：无缓存环几乎没有排队时延。"]),
            dict(t="ECN 标记", b=[
                "**代表**：S18 / S20。",
                "**听什么**：RSP 上的一比特标记。",
                "**优**：复用已有通道。**缺**：标记仍来自本节点失败。"]),
            dict(t="显式等级", accent=True, b=[
                "**代表**：S1 / S1T / S15 / S22。",
                "**听什么**：广播的拥塞等级或进度。",
                "**优**：全局可见。**缺**：等级同时含本核注入与被路过；S22 只换了线上内容。"]),
            dict(t="显式速率", b=[
                "**代表**：S28 / S28S。",
                "**听什么**：瓶颈自己算出的 share。",
                "**优**：源端不必猜。**缺**：总线最宽；不公平不在 hop 份额。"])])),

    ("matrix", dict(
        chrome="分类 · 控制点（第 21 页的列）",
        kicker="TAXONOMY BY CONTROL POINT  =  PAGE 19 COLS",
        title="动手处决定天花板（= 第 21 页的列）",
        cells=[
            dict(t="路径选择", b=[
                "**代表**：S26（UGAL / Valiant 类）。",
                "**动手处**：选最短还是绕远。",
                "**优**：不扣注入量。**缺**：绕远本身要多占链路。"]),
            dict(t="逐跳背压", b=[
                "**代表**：S27（PFC / wormhole 类）。",
                "**动手处**：上游链路 XOFF。",
                "**优**：不丢包，实现最简单。**缺**：需要可暂停队列，会形成拥塞树。"]),
            dict(t="源端速率", accent=True, b=[
                "**代表**：S1 族、S17 / S18 / S21 / S28。",
                "**动手处**：注入令牌桶。",
                "**优**：直接控总量。**缺**：闸门只能「扣住」，让出的槽被沿途吃掉。"]),
            dict(t="源端窗口", b=[
                "**代表**：S19 / S20。",
                "**动手处**：在飞事务数。",
                "**优**：与 CHI 事务模型对齐。**缺**：不管瓶颈处先服务谁。"]),
            dict(t="环上仲裁", b=[
                "**代表**：I-tag / S22。",
                "**动手处**：环上谁先用这一拍。",
                "**优**：控制点就在拥塞 hop。**缺**：要动仲裁器。"]),
            dict(t="预约调度", b=[
                "**代表**：S29（TDMA / Fastpass 类）。",
                "**动手处**：预分配时隙。",
                "**优**：拥塞不形成。**缺**：空闲时隙会浪费（S29 用 1 bit 回收）。"]),
            dict(t="接收端授权", b=[
                "**代表**：S16（Homa / NDP 类）。",
                "**动手处**：completer 的 DBIDResp。",
                "**优**：控制点在真正拥塞的下环侧。**缺**：要动事务层。"])])),

    ("figside", dict(
        chrome="S1 在分类中的位置", img="16a-cc-taxonomy.png", fig_w=8.30,
        caption="列 = 控制点 = 第 20 页（决定天花板）；"
                "行 = 拥塞信号 = 第 19 页（决定盲区）。"
                "灰底 = 已有方案，红框 = S1，绿框 = 本次补齐。",
        kicker="WHERE S1 SITS",
        blocks=[
            dict(kind="card", t="S1 = 源端速率 × 显式等级", b=[
                "**控制点**：「源端速率」列 —— 每核每窗的整型注入预算，出口是令牌桶。",
                "**信号**：「显式等级」行 —— 3 bit 拥塞等级，专用广播总线，固定 30 拍。"
                "同格还有 S1T 与 S15。"], wt=1.30),
            dict(kind="card", accent=True, t="这一格的信号语义", b=[
                "乘性减由**本节点上环失败**触发。无缓存环上这些失败多数来自他人 transit，"
                "所以等级同时含「本核注入」和「本核被路过」两件事。",
                "同一行里 S22 不改位宽，只把线上内容从「本窗有多堵」换成「本窗做了多少」，"
                "落后与领先可以直接比较。"],
                 wt=1.44),
            dict(kind="band",
                 text="62 组 AIMD 网格：带宽接近 S0 时 CoV 也回到 S0 附近，"
                      "没有同时明显改善两端的配置。", wt=0.74)])),

    ("matrix", dict(
        chrome="原本空着的四类", kicker="THE FOUR EMPTY CELLS",
        title="S0–S23 覆盖了六类，剩下四类本次补齐并实测",
        cells=[
            dict(t="S26 · 自适应路由", b=[
                "**原本为空**：全篇路由固定为「时延最短」，从未被重新计算过。",
                "**实现**：每核对两个出向 hop 各维护一个上环失败率 EWMA；"
                "最短方向比反向差 0.05 以上、且绕远 ≤ 2 跳，就改走反向。",
                "**代价**：无信号、无总线、无新报文，只多一个比较器 —— "
                "**1,560 FF‑eq**。"]),
            dict(t="S27 · 逐跳背压", b=[
                "**原本为空**：基线只有端到端的 CHI PCrd / tracker 信用，"
                "没有任何逐跳信号。",
                "**实现**：hop 占用率 ≥ 0.90 拉 XOFF、≤ 0.80 释放；"
                "沿相邻节点间的一根线逐跳上传，每跳 1 拍，reach = 2 跳。",
                "**代价**：每 (方向, VC) 一根相邻线，不是广播总线 —— "
                "**1,960 FF‑eq**。"]),
            dict(t="S28 · 显式速率反馈", accent=True, b=[
                "**原本为空**：S1 / S22 播的是拥塞等级或进度，没有方案播「速率」。",
                "**实现**：每个 hop 自己数本窗跨过它的核数 N 与占用率 y，"
                "按 RCP 的 share ← share×[1+α(C−y)/(C·N)] 更新，6 bit 播出；"
                "核取 min over 路径 hop 的 share ÷ 自己在该 hop 的流量占比。",
                "**代价**：40 个 hop 各一个 6 bit 字，总线最宽的一类 —— "
                "**21,480 FF‑eq**，全研究最贵。"]),
            dict(t="S29 · 预约 / 调度式", accent=True, b=[
                "**原本为空**：I-tag 是「饿了才预约」，没有方案做**事先**排班。",
                "**实现**：2 拍 × 10 核 = 20 拍帧，轮到某核时会骑过它出向 hop 的"
                "其他核让位（执行器与 S22 完全相同，只换触发）；"
                "每核 1 bit「我有 WriteData 排队」使空闲时隙不被浪费。",
                "**代价**：10 bit / 16 拍，比 S1 的总线窄得多 —— "
                "**4,440 FF‑eq**，是 S22 的 1/3。"])],
        band="选点口径：每一类的工作点都由该类自己的参数网格选出（probe_ring2_gapcc"
             "/ gapcc2），并取该类在三条轴上最好的那一行 —— "
             "所以一类如果输，是机制输，不是参数没调好。")),

    ("media", dict(
        chrome="补齐四类：机制与落点", kicker="FOUR NEW MECHANISMS",
        stat="4 / 8",
        stat_sub=["二维表里原本空着的格子数", "S26 / S27 / S28 / S29 已全部实现并实测"],
        img="16b-gap-diagram.png",
        caption="每一格给出该类在这条无缓存双向环上的落点：控制点、信号、"
                "以及执行器到底动了什么。四类都跑在与 S0 / S1 完全相同的 fabric 上。",
        cards=[
            dict(t="两类不用任何新线",
                 b=["S26 只看自己两个出向的上环失败率 EWMA；S29 只播 1 bit「我有 "
                    "WriteData 排队」。两者都不需要测量任何全局量。"]),
            dict(t="两类要加新信号",
                 b=["S27 每 (方向, VC) 一根相邻节点间的线，逐跳上传 reach = 2 跳；"
                    "S28 要 40 个 hop 各播一个 6 bit share，是全研究最宽的总线。"]),
            dict(t="口径与 S0 / S1 完全一致", accent=True,
                 b=[f"同一 K={n.kw}、同一 100 拍分箱、同一 outstanding = {n.oc}；"
                    "每类的工作点由该类自己的参数网格选出。"])])),

    ("figside", dict(
        chrome="补齐四类 vs S0 / S1", img="16c-gap-compare.png", fig_w=7.55,
        caption="同一 fabric、同一 K=20000、同一 100 拍分箱。三个数依次为"
                "总带宽 / 100 拍 CoV / 长期 max-min。",
        kicker="MEASURED, NOT ARGUED",
        blocks=[
            dict(kind="card", t="路由与背压：结构性失效", b=[
                f"**S26** {n.triple('S26')} —— 带宽 {n.dthr('S26')}，"
                "CoV 高于 S0：环已接近饱和，反向同样满，绕远还多占 (n−h) 个 hop·拍。",
                f"**S27** {n.triple('S27')} —— 带宽 {n.dthr('S27')}，"
                "不均衡度还比 S0 高。无缓存环上「保护链路」等于让链路空转。"], wt=1.06),
            dict(kind="card", t="显式速率：同一类的两个极点", b=[
                f"**S28**（RCP 反馈）{n.triple('S28')} —— 带宽 {n.dthr('S28')}，"
                "公平几乎不动："
                "不公平不在 hop 的份额分配，在每核自己上环口的仲裁，"
                "算出的 share 高于落后核能上环的速率。",
                f"**S28S**（每 hop 静态等分）{n.triple('S28S')} —— "
                f"全研究最低的不均衡度，代价带宽 {n.dthr('S28S')}。"], wt=1.12),
            dict(kind="card", accent=True, t="S29 与 S1 的三轴数字", b=[
                f"S29：{n.triple('S29')}；S1：{n.triple('S1')}。"
                "总线 10 bit / 16 拍，对 S1 的 6 bit + 20 项表。",
                "S29 复用 S22 的让位：让出的槽是指名的。"
                "官方 K=20000 工作点见后文实测页。"], wt=1.30)])),

    ("section", dict(
        no="05", kicker="SECTION FIVE", title=["公平性 — 总带宽", "的 trade-off"],
        lead="先把「要少一点不均衡到底该付多少带宽」这条曲线精确算出来，"
             "再从它推导出一个同时表示总带宽和不均衡度的标量 φ。",
        key_label="KEY MESSAGE",
        key=f"总带宽和公平性确实不可兼得：把十个核压到完全等速率，"
            f"结构上就要放弃 {(1 - n.rstar/n.rmax)*100:.2f}% 的峰值带宽。"
            f"理想汇率 κ = {n.kappa:.3f} flit/cycle 每单位 CoV；"
            f"φ = (R − κ·CoV)/R* 在上界上恒为 1，"
            f"S16 {n.phi('S16')}、S0 {n.phi('S0')}、S1 {n.phi('S1')}。")),

    ("figside", dict(
        chrome="公平性 — 带宽交换曲线", img="16-tradeoff.png", fig_w=8.15,
        caption="横轴 = 不均衡度 CoV（十核 100 拍窗带宽的标准差 / 均值，0 = 完全均等，越右越不均）；"
                "红线 = 每个 CoV 上限下的理论最高带宽；"
                "13 个点 = 官方 K = 20000 全部实测方案（含 S26–S29）。"
                "I-tag 只是 S0 的 t_inj / hold 调参，不是独立机制。",
        kicker="HOW TO READ THE CURVE",
        blocks=[
            dict(kind="card", t="先认四个量", b=[
                "**λc**：第 c 个核的事务速率。",
                "**占用表**：每笔事务占哪些链路、各几拍；路由固定后这张表就固定。",
                "**CoV**：允许的最大不均衡度（十核速率的标准差 / 均值）。",
                "**R(CoV)**：不均衡度 ≤ CoV 时的最高总带宽。"],
                 wt=1.70),
            dict(kind="band",
                 text="红线点：固定 CoV 上限，枚举十个核的速率组合，链路每拍 ≤ 1 flit，"
                      "取总带宽最高的那个。CoV 从 0 逐点放宽即得整条红线"
                      "（凸优化全局上界，不是对仿真点拟合）。", wt=1.10),
            dict(kind="band",
                 text="红线上 = 该不均衡度下理论最好；红线下 = 机制损失；红线上方不可能。"
                      "最左端 CoV = 0 即 R* = 5.7143；CoV ≥ 0.306 后红线平在 R_max = 6.40。"
                      "在 CoV 坐标下红线是直线 R* + κ·CoV，κ = 2.18 —— 为什么见后两页。", wt=0.90)])),

    ("figside", dict(
        chrome="标量指标推导：把上界在 CoV 坐标下拉直", img="42-metric-derivation.png",
        fig_w=7.05,
        caption="左：LP 上界的 79 个点画在（CoV，总带宽）平面上是一条直线 R = R* + κ·CoV"
                "（κ = 2.177，RMS 残差 0.011，最大相对残差 0.43%）；蓝虚线 = 过 S0 的平行线，"
                "它与纵轴的交点就是 Φ(S0)。右：13 个方案 1 − φ 的两项分解。",
        kicker="SYMBOLS, THEN HOW TO READ IT",
        blocks=[
            dict(kind="card", t="符号表（单位都是带宽 flit/cycle）", b=[
                "**R** 实测总带宽；**R*** = 5.7143，十核完全等速率（CoV = 0）时的最高带宽。",
                "**CoV** 不均衡度 = 十核 100 拍窗带宽的标准差 / 均值。",
                "**κ** = 2.177，红线斜率 = 理想汇率：最好的控制器每多 1 单位 CoV 换到的带宽。",
                "**Φ = R − κ·CoV**：把不均衡按理想汇率「退」成带宽后剩下的等速率等效带宽；"
                "**φ = Φ / R***，理想 = 1。"],
                 wt=1.70),
            dict(kind="card", accent=True, t="左图怎么读", b=[
                "红点全在红线上 ⇒ 上界任一点 φ ≡ 1。过 S0 作红线的平行线（蓝虚线），"
                f"滑到 CoV = 0 读出 **Φ(S0) = {n.W('S0')['throughput'] - n.kappa * n.covf('S0'):.3f}**；"
                f"竖箭头 R* − Φ = **{(1 - n.phif('S0')) * n.rstar:.3f}** = (1−φ)·R*。",
                "灰虚线 = 等 φ 线，与红线平行：线上每一点的等速率等效带宽 Φ 相同。"],
                 wt=1.40),
            dict(kind="band",
                 text="右图：1 − φ = 蓝（带宽缺口 (R*−R)/R*）+ 橙（不均衡折成带宽 κ·CoV/R*）。"
                      f"S16 {n.phi('S16')}、S0 {n.phi('S0')}、S1 {n.phi('S1')}。",
                      wt=0.58)])),

    ("triple", dict(
        chrome="为什么理想控制器在带宽—CoV 图中是一条直线", kicker="WHY A STRAIGHT LINE",
        title="最优的「不公平」永远是同一个方向，只是幅度不同",
        cards=[
            dict(num="01", t="理想控制器只有一根杠杆",
                 b=["限制总带宽的是热段上的 4 条绑定链路。邻 mem = 1 的 4 个慢核"
                    "（C0 / C8 / C10 / C18）的写只能从一侧出去，大部分压在绑定链路上；"
                    "其余 6 核可以两侧分流，每笔写只有一小部分落在绑定链路。",
                    "所以「以不均衡换带宽」只有一种划得来的做法：把速率从 4 个慢核挪给 6 个快核。"
                    "每个 CoV 上限下 LP 的最优解都是这同一个方向，只是挪的量 δ 不同。"]),
            dict(num="02", t="每挪一单位，带宽和标准差各多一个常数", accent=True,
                 b=["挪 δ：绑定链路省出的容量让总带宽多 **12/7 ≈ 1.71·δ**；"
                    "十核速率的标准差多 **√24/10 ≈ 0.49·δ**（4 对 6 分组的几何常数）。"
                    "两者都正比于 δ ⇒ 带宽对标准差是精确的直线。",
                    "CoV = 标准差 / 均值，均值全程只从 0.571 变到 0.640（+12%），"
                    "所以对 CoV 也几乎是直线：精确式 **R = R* / (1 − 0.35·CoV)**，"
                    "在 0–0.306 内与直线最大偏差 0.43%。"]),
            dict(num="03", t="这条直线代表什么",
                 b=["**斜率 κ = 2.18 是汇率**：最好的控制器用 0.1 CoV 的不均衡恰好换到 0.22 flit/cycle，"
                    "在线上哪一点都一样。",
                    "**左端** CoV = 0：max-min 公平，R* = 5.714。**右端** CoV = 0.306：吞吐最大 R_max = 6.40，"
                    "慢核速率只有快核的一半；再不均也换不到带宽，线变平。",
                    "**线上的点都「不浪费」**：少掉的每份带宽都按理想汇率换成了均衡；"
                    "线下的点到线的竖直距离 = 同等 CoV 下损失的带宽。"])],
        band="因为方向固定、汇率恒定，任何方案都可以沿这个方向投影回 CoV = 0 —— 投影落点的带宽就是 Φ，"
             "除以 R* 就是 φ。这就是把「总带宽 + 不均衡度」两个数压成一个数而不丢信息的依据。")),

    ("triple", dict(
        chrome="其他方案在带宽—CoV 图中是什么形状", kicker="WHAT REAL SCHEMES LOOK LIKE",
        title="一个工作点是一个点；调参数会走出三种轨迹",
        cards=[
            dict(num="01", t="工作点：全部落在线下方",
                 b=[f"13 个方案各是一个点（第 27 页左图），没有一个在线上或线上方。"
                    f"S0 ({n.cov('S0')}, {n.thr('S0', 3)})：带宽 {n.pct_r('S0')}% R*，"
                    f"但多付 {n.cov('S0')} 的不均衡；"
                    f"S28S ({n.cov('S28S')}, {n.thr('S28S', 3)})：够均衡但带宽 {n.dthr('S28S')}；"
                    f"S27 在右下角：两头都差。",
                    f"离红线越远 φ 越低：S16 {n.phi('S16')}、S0 {n.phi('S0')}、"
                    f"S1 {n.phi('S1')}、S27 {n.phi('S27')}。"]),
            dict(num="02", t="调参轨迹：三种走法", accent=True,
                 b=["**沿等 φ 线**（斜率 κ）：与理想同效率地以带宽换均衡 —— S16 oc 16 → 12。",
                    "**竖直向下**（CoV 不变、只掉带宽）：S16 oc < 12、S29 slot ≥ 3、S1 harsh 档。",
                    "**中间斜线**（斜率 κ′ > κ，ε = κ/κ′ < 1）：S1 cap 段 ε 0.68–0.79，"
                    "S29 slot 1 → 2 0.82，S28 0.31–0.80。",
                    "过拐点后折向右下：S1 spec → harsh 带宽掉、CoV 反而回升。"]),
            dict(num="03", t="为什么是这些形状",
                 b=["走等 φ 线要「只从快核拿、只给慢核」且链路不空转 —— "
                    "S16 在 HA 授权点重排服务顺序、槽位一个不丢，恰好如此。",
                    "竖直下降 = 对十核一律限流：份额比例不变（CoV 不变）、链路空转（带宽掉）。"
                    "全局阈值一旦压到慢核本来就到不了的速率，就进入这一段。",
                    "中间斜线 = 两者混合：S1 的拥塞等级把热段附近的慢核也压了；"
                    "S27 背压让链路空转又不重排份额，所以落到右下角。"])],
        band="第 31 页把 13 个方案各自的旋钮轨迹画全：只有 S16 的一段贴近等 φ 线，"
             "其余都是折线、竖线或一团点 —— 旋钮不是 LP 那根杠杆。")),

    ("matrix", dict(
        chrome="φ 怎么用：表征一种算法、预测不均衡度、比较方案",
        kicker="WHAT φ, Δφ AND φ-RATIOS MEAN",
        title="调参降带宽 → 用 φ 反推 CoV；两方案相减 → 同等均衡下的带宽差",
        cells=[
            dict(t="① φ 表征一种算法", b=[
                "一个方案的 (R, CoV) 工作点给出 φ。调参时沿**等 φ 线**移动的算法，"
                "就是「以理想汇率把带宽换成均衡」的算法。",
                "φ 不变、带宽降到 R′ 时：**CoV′ = (R′ − φ·R*)/κ**。"
                "预测值是下界：实测 CoV 不会低于它。"]),
            dict(t="② 预测的验证", b=[
                "筛选轮（K=2000）上：S16 oc 16→12、S29 slot 2→1、S1 spec cap "
                "沿等 φ 线移动时，预测 CoV 是实测的下界。",
                "过拐点的档（S1 spec→harsh）预测失效 —— 见 ③。",
                f"官方 K={n.kw} 工作点用同一公式现算：S16 φ = {n.phi('S16')}，"
                f"S0 φ = {n.phi('S0')}。"]),
            dict(t="③ 拐点与交换效率 ε", b=[
                "每段旋钮的实测斜率 κ′ = ΔR/ΔCoV，**ε = κ/κ′** = 实际换到的均衡占理想汇率的比例。",
                "筛选轮上 S16 的 oc 段 ε 最高；S29 slot、S1 cap、S28 都明显低于 1。",
                "过拐点后 ΔCoV ≥ 0、ε ≤ 0：只掉带宽不换均衡，φ 直接按 ΔR/R* 下降。"]),
            dict(t="④ 差值与比值的含义", accent=True, b=[
                "**Δφ·R*** = 同等不均衡度下 A 比 B 多出的带宽（把两者的 CoV 差按理想汇率折成带宽后相减）。"
                f"S16 − S0 = **{n.phif('S16')-n.phif('S0'):+.3f}**；"
                f"S1 − S0 = **{n.phif('S1')-n.phif('S0'):+.3f}**。",
                f"**φ_A / φ_B** = 等速率等效带宽之比："
                f"S16/S0 = {n.phif('S16')/n.phif('S0'):.3f}，"
                f"S0/S1 = {n.phif('S0')/n.phif('S1'):.3f}。",
                "**(1−φ) 之比** = 离理想控制器的距离之比。"])],
        band="φ 的差值可直接读成 flit/cycle，并按理想汇率拆成带宽项与均衡项。")),

    ("figside", dict(
        chrome="13 个方案各自的旋钮轨迹：都不是那根 LP 直线",
        img="43-metric-knobs.png", fig_w=7.55,
        caption="每个面板 = 该方案自己的一个旋钮；K = 2000，CoV 仍是 100 拍窗。"
                "红线 = LP 上界（CoV ≥ 0.306 后平在 R_max）；星 = 官方工作点。"
                "S1 按 band 拆成三条，不是把 9 个点连成一条。",
        kicker="KNOB ≠ LP LEVER",
        blocks=[
            dict(kind="card", t="是的：每条曲线是该方案自己的旋钮", b=[
                "每个面板只动一个标量：S0 / ITAG 扫 t_inj，S16 扫 oc，"
                "S1 按 band 拆成三条，S29 扫 slot。",
                "红星 = 官方点。13 次独立扫描，不是混拟合。"],
                 wt=1.00),
            dict(kind="card", accent=True, t="为什么大部分不是直线", b=[
                "红线直，因为最优只有一个方向：4 个慢核 → 6 个快核，"
                "R 与 σ 都正比于挪动量 δ。",
                "旋钮不是这个 δ：先重分配，再一律限流（竖着掉），"
                "再错杀 / 空转 / 短窗抖动。"],
                 wt=1.10),
            dict(kind="card", t="本质：旋钮 ≠ LP 的那根杠杆", b=[
                "LP 直线是最优速率沿单一自由度的像；"
                "旋钮同时搅动份额、空转和 100 拍窗抖动。",
                "筛选轮上仅 S16 的 oc 段贴近红线斜率。S28S 竖直；"
                "S19 / S20 / S22 / S1T 不动；S26 两边都坏。"],
                 wt=1.10)])),

    ("section", dict(
        no="06", kicker="SECTION SIX", title=["理论最优与", "各方案 Pareto"],
        lead="先钉死「无限聪明、无限快的控制器最多能做到什么」，"
             "再把实测方案摆进同一张图。",
        key_label="KEY MESSAGE",
        key=f"一个无限聪明的拥塞控制器，最好也就是 {n.rstar:.4f} flit/cycle"
            f"（S0 的 {100.0 * n.rstar / n.W('S0')['throughput']:.1f}%）"
            "配 CoV ≈ 0。可争的总带宽只有约 3 个点，均衡那一侧才是主战场。")),

    ("triple", dict(
        chrome="理想拥塞控制器的上限",
        kicker="THE REFERENCE EVERYTHING IS DIVIDED BY",
        title="无限聪明、无限快，只受资源守恒约束",
        cards=[
            dict(t="什么是 LP 模型", b=[
                "LP = 线性规划：**变量**是十个核各自的注入速率，"
                "**约束**是每条链路每拍最多搬 1 flit，**目标**是把总带宽做到最大。"
                "因为路由和握手都固定，「一笔事务占几条链路各几拍」是一张常数表，"
                "占用量随速率**线性**增长，所以这个最大值能被精确解出来、不需要仿真。",
                "模型里**没有**缓冲、没有在环优先、没有 I-tag / E-tag —— "
                "这是有意的：让机制造成的损失也留在待解释的差距里。"]),
            dict(t="两个界", accent=True, b=[
                "**最大吞吐** R_max = 6.4000 flit/cycle。但这一点的解**不唯一**，"
                "其中最差的那个把 C0 / C10 完全饿死 —— 不是可交付的工作点。",
                "**等速率点** = 强行要求十个核速率**完全相同**时能达到的最大总带宽。"
                "此时每核 λ* = 2/7 = 0.2857 笔/拍，合计 R* = 40/7 = **5.7143**。"
                "它就是「最公平」那一端的天花板，全篇的分母都是它。",
                "瓶颈在环上的链路，不在注入或下环端口 —— 加深队列、加宽端口都抬不高它。"]),
            dict(t="理想控制器能做到什么", b=[
                "100 拍箱内全环有 571.4 个写 flit，每核约 57.1 个。"
                "一个确定性的理想控制器把箱内总量按整数尽量均分，得 "
                "CoV_ideal = **0.0055**。",
                "也就是说：**理论最优 = 5.7143 flit/cycle（S0 的 103.1%）配 CoV ≈ 0**。"
                "整数粒度不构成任何限制。"])],
        band="适用范围：λ* = 2/7 是「fabric + 流量 pattern」的联合解，不是 fabric 常数。"
             "换成 hot 之后绑定资源从入环 hop 变成热簇的下环口，λ* 掉到 0.100、"
             "R* 掉到 2.0000 —— 任何把 λ* 写进硬件的设计只对它被推导时的那个 pattern 正确。")),

    ("figside", dict(
        chrome="收益 — 硬件开销 Pareto", img="18-pareto.png", fig_w=7.60,
        caption="纵轴 φ = (R − κ·CoV)/R*：总带宽减去按理想汇率折成带宽的不均衡，再除以 R*，理想 = 1.0；"
                "筛选轮 K = 2000，100 拍窗。横轴 FF‑eq = 等效触发器数，"
                "把新增的寄存器、比较器、加法器、乘法器与总线线宽都折算成"
                "「相当于多少个 D 触发器」。",
        kicker="FRONTIER",
        blocks=[
            dict(kind="card", t="硬件 Pareto 来自 K=2000 筛选轮", b=[
                "左图是筛选轮的面积—φ，用来排除明显不值的点；"
                f"官方口径是 K={n.kw}、outstanding = {n.oc}，φ 用下一句现算。",
                f"官方点：S0 φ = {n.phi('S0')}（0 FF‑eq）→ "
                f"[[S16]] φ = **{n.phi('S16')}**（900 FF‑eq）；"
                f"S22 φ = {n.phi('S22')}（13,920）；S29 φ = {n.phi('S29')}（4,440）；"
                f"S1 φ = {n.phi('S1')}（21,220）。"],
                 wt=1.20),
            dict(kind="card", t="补齐的四类：官方点仍未进前沿", b=[
                f"S26 φ = {n.phi('S26')}；S27 φ = {n.phi('S27')}；"
                f"S28 φ = {n.phi('S28')}；S28S φ = {n.phi('S28S')}；"
                f"S29 φ = {n.phi('S29')}。",
                f"除 S29 外，四个新点的 φ 都低于 S0 的 {n.phi('S0')} —— "
                "不是被 S16 挤掉，是连「什么都不做」都没超过。",
                "S29 是环仲裁分支上最便宜的点（4,440 FF‑eq，S22 的 1/3）。"],
                 wt=1.42),
            dict(kind="band",
                 text=f"判据是「离理想控制器有多近」：官方 K={n.kw} 上 "
                      f"S16 把 φ 从 {n.phi('S0')} 推到 {n.phi('S16')}。",
                 wt=0.62)])),

    ("figside", dict(
        chrome="固定非均匀流量下的带宽 Pareto", img="20-hot-pareto.png", fig_w=7.60,
        caption="hot：十个核全部写入 HA 11 / 13 两节点簇。R* 用它自己的目的地分布重解 "
                "= 2.0000，绑定资源变成热簇的下环口。此处**只看总带宽**。",
        kicker="ROBUSTNESS TO TRAFFIC SHIFT",
        blocks=[
            dict(kind="card", accent=True, t="为什么这里不看瞬时 CoV", b=[
                "十个核全部打同一个两节点热簇，谁快谁慢由「离热簇几跳」直接决定，"
                "拥塞点也不在注入侧 —— 这一页只考察**总带宽是否被换掉**。"], wt=0.86),
            dict(kind="card", t="同一 outstanding = 128 时，hot 先打容量", b=[
                "十核全打 HA 11/13，在飞窗口不再被 RTT 限住，"
                "S0 只拿到 R* 的 **74.0%**。窗口类（S19 / S20）收到 **98.9%**，"
                "因为它们会把 requester 窗口从 128 收紧。",
                "[[uniform 写下 outstanding = 128 不绑定；hot 上它就是过量注入。]]"],
                 wt=1.20),
            dict(kind="card", t="控制点仍然决定谁能救带宽", b=[
                "S16 把决策放在目的地，但 overcommit 是为 uniform 占用调的；"
                "hot 上 completer 更满，同一 16 会过狠（筛选轮 57.8% R*）。",
                "S28 / S28S 显式速率保住 94–97%；S27 背压掉到 72%。"
                "换 pattern 之后，**先保容量的是窗口，不是授权重排**。"],
                 wt=1.40),
            dict(kind="band",
                 text="hot 这一页回答的是「换流量会不会把带宽换掉」，"
                      "不是把 uniform 的公平结论原样搬过来。", wt=0.54)])),

    ("figside", dict(
        chrome="各方案的各核完成时间曲线", img="44-finish-all.png", fig_w=8.60,
        caption="13 个方案，同一 fabric、K = 20000 uniform 写；每格十条线 = 十核累计上环 "
                "WriteData / 配额，虚线 = 最早与最晚完成时刻。S1U / S1D 与 S1 / S0 逐拍相同，不重复。",
        kicker="WHO FINISHES LAST",
        blocks=[
            dict(kind="card", t="三种形态", b=[
                f"**六快四慢**（S0 {n.frat('S0')}、S1T {n.frat('S1T')}、"
                f"S19 {n.frat('S19')}、S20 {n.frat('S20')}）："
                f"S0 完成时间 {n.fspan('S0')}。",
                f"**收得很齐**（S16 {n.frat('S16')}、S28S {n.frat('S28S')}）。",
                f"**中间态**（S22 {n.frat('S22')}、S29 {n.frat('S29')}、"
                f"S1 {n.frat('S1')}）：慢核仍慢，差距收窄。"],
                 wt=1.45),
            dict(kind="card", accent=True, t="拐点的含义", b=[
                "六快四慢的方案里，慢核曲线在快核完成瞬间斜率抬升 —— "
                "瓶颈是满载环上的上环机会，不是慢核自身。",
                f"S28S 够齐（比 {n.frat('S28S')}）但总时长被拉到 "
                f"{n.finish_span('S28S')[1]/1000:.1f}k，带宽 {n.dthr('S28S')}。"],
                 wt=1.15),
            dict(kind="band",
                 text="完成时间比是长期速率比的另一种读法；它和 100 拍 CoV 不是同一件事，见下页。",
                 wt=0.45)])),

    ("media", dict(
        chrome="完成时间跨度 vs 瞬时不均衡度",
        kicker="TWO READINGS",
        stat=f"{n.finish_span('S16')[1] - n.finish_span('S16')[0]:,} 拍",
        stat_sub=[f"S16 十核最早与最晚完成时刻之差（{n.fspan('S16')}）",
                  f"S0 为 {n.finish_span('S0')[1] - n.finish_span('S0')[0]:,} 拍"],
        img="45-finish-spread.png",
        caption="左：每个方案十核从最早到最晚完成的区间（千拍），右端数字 = 最晚 / 最早；"
                "右：100 拍 CoV 与完成时间比的散点。",
        cards=[
            dict(t="两个口径给出不同排序",
                 b=[f"S27 CoV {n.cov('S27')} 高于 S1 的 {n.cov('S1')}，"
                    f"完成时间比却是 {n.frat('S27')} 对 {n.frat('S1')}："
                    "背压把所有核一起拖慢，长期份额反而接近。"
                    f"S22 CoV {n.cov('S22')}，完成时间比 {n.frat('S22')}。"]),
            dict(t="S16 两项一起动", accent=True,
                 b=[f"S16：CoV **{n.cov('S16')}**、完成时间比 **{n.frat('S16')}**、"
                    f"总带宽 {n.dthr('S16')}。"
                    f"S28S 完成时间比 {n.frat('S28S')}，但带宽 {n.dthr('S28S')}。"]),
            dict(t="对架构评审的意义",
                 b=["瞬时 CoV 回答「任一 100 拍里有没有核被饿」，完成时间比回答"
                    f"「一批工作最后是谁拖尾」；两者都要看，S1 的 CoV {n.dcov('S1')} "
                    f"对应完成时间比 {n.frat('S1')}（S0 为 {n.frat('S0')}）。"])])),

    ("section", dict(
        no="07", kicker="SECTION SEVEN", title=["前沿方案详述"],
        lead="三个落地候选、两个研究对照：机制示意 + 与 S0 / S1 的同口径实测；"
             "S16 / S22 / S29 各再加两页微架构框图与逐拍示例。",
        key_label="SELECTION RULE",
        key="S16 是主方案；事务层不可改时走环仲裁分支，其中 S22 效果最好、"
            "S29 便宜 3.1 倍；S19 / S20 只用于验证 requester 动态窗口的边界。")),

    ("figside", dict(
        chrome="S16 授权保留（900 FF‑eq） · receiver-driven / 本地触发 / 零新报文",
        img="22-s16-diagram.png", fig_w=8.75,
        caption="上 = 写路径：把协议本来就要发的 DBIDResp 当授权用；"
                "下 = 读路径：同一套状态可以照搬，但实测不需要（见下页）。",
        kicker="FRONTIER · COMPLETER MECHANISM",
        blocks=[
            dict(kind="card", t="四段机制（与 S1 同格式）", b=[
                "**① 检测**：HA 只数自己的在飞授权数。不需要任何拥塞信号 —— "
                "拥塞就发生在它自己的下环口上。",
                "**② 传递**：零。不加总线、不加带内标记，线格式完全不动。",
                f"**③ 反馈**：REQ 到达后先排队而不是立即授权，"
                f"同时在飞授权 ≤ overcommit = {n.oc16}。",
                "**④ 执行**：在排队者里把授权给**迄今被服务最少**的那个核。"],
                 wt=1.62),
            dict(kind="card", accent=True, t="改的是事务层，不是链路层", b=[
                "四段全部落在 completer 内部：它决定**什么时候、先给谁**发 DBIDResp。"
                "线格式、协议状态机都不动，但**动了 HA 的事务调度** —— "
                "这一条决定它能不能落地，见结论页。"], wt=1.06),
            dict(kind="band",
                 text=f"overcommit = {n.oc16} 跟踪的是 HA 实际占用（≈ RTT×λ / n_HA），"
                      f"不是 core_outstanding：上限 {n.oc} 覆盖最坏 RTT 之后占用仍约 30，"
                      f"工作点仍是 {n.oc16}。取到 {n.oc} 就扣不住授权，S16 退化成 S0。",
                 wt=0.72)])),

    ("media", dict(
        chrome="S16 实测：写侧要做，读侧不用做", kicker="MEASURED · UNIFORM",
        stat=n.cov("S16"),
        stat_sub=[f"S16 写侧 100 拍 CoV（S0 = {n.cov('S0')}）",
                  f"总写带宽 {n.dthr('S16')}，整窗 max/min = {n.mm('S16')}"],
        img="23-s16-compare.png",
        caption="左两幅 = 写（K = 20000），右两幅 = 读（K = 5000）；"
                "每幅内三根柱依次 S0 / S1 / S16，橙虚线 = 理论上限 R*。"
                "注意读侧 S0 与 S16-R 两根柱几乎一样高。",
        cards=[
            dict(t="写：低代价把份额搬回去", accent=True,
                 b=[f"CoV {n.cov('S0')} → **{n.cov('S16')}**，整窗 max/min "
                    f"{n.mm('S0')} → **{n.mm('S16')}**，"
                    f"总带宽 {n.dthr('S16')}，φ {n.phi('S0')} → {n.phi('S16')}。",
                    f"对照 S1：CoV {n.cov('S1')}，带宽 {n.dthr('S1')}。"]),
            dict(t="读：S0 本来就是齐的（128 B CompData）",
                 b=[f"十个核读侧 CoV {n.rcov('S0')}、max/min {n.rmm('S0')}，"
                    f"带宽 {n.rpct('S0')}% R*。"
                    "**128 B 读没有待解决的问题**；64 B 读见后两页。"]),
            dict(t="所以读侧建议不做",
                 b=[f"S16-R 带宽 {n.rpct('S16-R')}% R*（{n.rdthr('S16-R')}）、"
                    f"CoV {n.rcov('S16-R')}，"
                    "[[收益在噪声量级，不值一次事务层改动。]]",
                    f"对照：S1-R 读侧带宽 {n.rdthr('S1-R')}、CoV {n.rcov('S1-R')}。"])])),

    ("media", dict(
        chrome="读侧补充：CompData = 1 / 2 / 4 flit 时 S0 的不均衡度",
        kicker="READ-SIDE FAIRNESS",
        stat=n.pmm("S0-m1"),
        stat_sub=[f"CompData = 1 flit 时 S0 读侧 max/min（2 flit = {n.pmm('S0-m2')}）",
                  f"与写侧的 {n.mm('S0')} 同量级，同一组慢核"],
        img="40-read-payload.png",
        caption="S0 读，K = 5000 笔/核，REQ 1 flit，CompData 1 / 2 / 4 flit（64 / 128 / 256 B）；"
                "右一横轴归一到各自的争用窗长度。",
        cards=[
            dict(t="1 flit：不均最大", accent=True,
                 b=[f"max/min **{n.pmm('S0-m1')}**，100 拍 CoV **{n.pcov('S0-m1')}**；"
                    "C0 / C10 最慢、C8 / C18 次之 —— 与写侧同一组慢核。"
                    f"总带宽 {n.pthr('S0-m1')}。"]),
            dict(t="2 flit 齐；4 flit 长期齐、瞬时抖",
                 b=[f"2 flit：max/min {n.pmm('S0-m2')}、CoV {n.pcov('S0-m2')}。",
                    f"4 flit：max/min **{n.pmm('S0-m4')}**，CoV **{n.pcov('S0-m4')}** ——"
                    "每笔 4 flit 的突发让 100 拍箱内份额抖动，长期份额不受影响。"]),
            dict(t="为什么 1 flit 会不均",
                 b=["1 flit 时 REQ（core→HA）与 CompData（HA→core）每笔各 1 flit，两向负载相同，"
                    "REQ VC 与 DAT VC 并列最忙（占用 0.943），写侧的注入几何差异经 REQ 方向回到读侧。",
                    "2 / 4 flit 时 DAT 单独成为瓶颈（占用 0.980 / 0.964），CompData 由 HA 侧发出，各核等齐。"])])),

    ("media", dict(
        chrome="读侧各核完成时间：1 flit 六快四慢，2 / 4 flit 十线重合",
        kicker="READ COMPLETION",
        stat=f"{n.pfrat('S0-m1')} / {n.pfrat('S0-m2')} / {n.pfrat('S0-m4')}",
        stat_sub=["十核完成时间 最晚 / 最早：CompData = 1 / 2 / 4 flit",
                  "只有 1 flit 出现写侧的形态"],
        img="41-read-payload-finish.png",
        caption="同上三组；纵轴 = 已收到 CompData / 本核配额，红 = 邻 mem = 1 的四核，蓝 = 其余六核。",
        cards=[
            dict(t="1 flit", accent=True,
                 b=[f"完成时间 {n.pfspan('S0-m1')}，比值 {n.pfrat('S0-m1')}；"
                    "慢核曲线在快核完成后抬升，与写侧 S0 同一形态。"]),
            dict(t="2 / 4 flit",
                 b=[f"{n.pfrat('S0-m2')} 与 {n.pfrat('S0-m4')}，十条线不可分。"
                    "4 flit 的抖动只在 100 拍尺度出现，累计曲线上看不到。"]),
            dict(t="对建议的影响",
                 b=["若 workload 含大量 64 B 读，读侧需要和写侧一样纳入 S16-R；"
                    "128 B 及以上读维持现状的结论不变。"])])),

    ("figside", dict(
        chrome="S16 微架构：HA 内部新增的三张小表和一条决策规则",
        img="32-s16-uarch.png", fig_w=8.15,
        caption="灰框 = 基线 HA 已有；红框 = S16 新增。REQ 从左侧下环进 tracker，"
                "授权决策读三张小表后驱动已有的 DBIDResp 发送器。",
        kicker="MICROARCHITECTURE · S16",
        blocks=[
            dict(kind="card", t="新增状态（每 HA 一份）", b=[
                f"**outstanding**：{max(5, (n.oc16).bit_length())} bit，在飞授权数 "
                f"（上限 {n.oc16}）；授权 +1，一笔写的末 flit 落地 −1。",
                "**served[c]**：10 × 10 bit 饱和计数，每授权 +2 flit，定期同减最小值。",
                "**待授权计数**：10 × 6 bit，外加 tracker 每条目 1 bit「已授权」标志。"],
                 wt=1.45),
            dict(kind="card", accent=True, t="决策只有两条规则", b=[
                f"直通：在飞 < {n.oc16} 且无人等待 → 到达即授权，轻载零延迟。",
                "排队：从有等待的核里选 served 最小者（同值取小核号），"
                "取其最老的未授权条目发 DBIDResp。"], wt=1.00),
            dict(kind="band",
                 text="无周期控制环、无总线、无新报文。成本口径 900 FF‑eq："
                      "10 bit 计数 + 2 比较 + 1 加法，×10；表与标志位挂在 tracker 上未单列。",
                 wt=0.70)])),

    ("figside", dict(
        chrome="S16 工作示例：授权是怎么流向落后核的",
        img="33-s16-flow.png", fig_w=8.15,
        caption=f"示意 overcommit = 2（实际 {n.oc16}，逻辑逐字相同）。左 = 报文时序；"
                "右 = HA 三张表在每一步之后的值，红字是那次改变结果的比较。",
        kicker="WALKTHROUGH · S16",
        blocks=[
            dict(kind="card", t="六步看懂", b=[
                "①② C0 两笔 REQ 直通授权，在飞到 2。③ C2 的 REQ 只能等。④ C0 第三笔也等。",
                "⑤ C0 第一笔写完 → 空出的名额给 served 最小的 C2，不是先到的 C0。",
                "⑥ 再空一个名额才轮到 C0。"], wt=1.25),
            dict(kind="card", accent=True, t="为什么它几乎不花带宽", b=[
                "名额一空立刻补，在飞授权始终顶在上限，HA 下环口没有一拍空转。",
                f"它改的是**谁**拿授权，不是**多少**授权 —— "
                f"所以总带宽 {n.dthr('S16')}。"],
                 wt=1.00),
            dict(kind="band",
                 text="读侧同一套状态可照搬（served 计 CompData，末 CompData 落地释放名额），"
                      "但前页已证明读侧没有待解决的问题。", wt=0.72)])),

    ("figside", dict(
        chrome="S19 Swift / S20 DCTCP（各 5,840 FF‑eq） · requester 动态窗口",
        img="28-window-diagram.png", fig_w=8.75,
        caption="两者共用每核动态 outstanding 窗口；S19 用协议天然 RTT，"
                "S20 用 HA tracker 产生的 1 bit ECN 标记。",
        kicker="REFERENCE · WINDOW CC",
        blocks=[
            dict(kind="card", t="共同执行器：动态 outstanding", b=[
                "每个 core 维护窗口 Wc：只有当前在飞事务数 < Wc，才允许发新的 REQ。"
                f"初值 16、下限 8、硬上限仍是 core_outstanding = {n.oc}。",
                "事务完成归还名额；Retry 重发不受窗口限制。"
                "窗口有名额时允许突发，因此比逐拍 rate pacer 更贴近 CHI 的事务模型。"],
                 wt=1.25),
            dict(kind="card", accent=True, t="S19：RTT 驱动", b=[
                "REQ 上环到 DBIDResp 回到 core 就是一份现成 RTT 样本，"
                "Retry 往返也被计入。RTT 低于 target 就加窗，高于 target 就按超额比例缩窗。"
                "**不加报文、不加标记位。**"], wt=1.00),
            dict(kind="card", t="S20：ECN 驱动", b=[
                "HA 根据 request tracker 占用率在 DBIDResp 上附 1 bit mark，RetryAck 直接视为 mark。"
                "core 对标记比例做 EWMA；有标记按 α/2 缩窗，无标记按 1/W 加窗。"], wt=1.00)])),

    ("media", dict(
        chrome="S19 / S20 实测：均匀写", kicker="MEASURED · UNIFORM",
        stat=n.dthr("S19"),
        stat_sub=["S19 相对 S0 的总写带宽（S20 见右卡）",
                  "两者各 5,840 FF-eq · requester 动态窗口"],
        img="29-window-compare.png",
        caption="三幅依次比较总带宽、100 拍瞬时 CoV、整窗最快 / 最慢核带宽比；"
                "每幅内四根柱都是 S0 / S1 / S19 / S20。",
        cards=[
            dict(t="信号不同，执行器相同", accent=True,
                 b=["S19 看端到端 RTT，能覆盖 ring 与 completer 等待；"
                    "S20 只看 HA tracker 压力，信号更直接但需要 DBIDResp 的 1 bit mark。"]),
            dict(t="当前工作点结果接近",
                 b=[f"**S19**：带宽 {n.thr('S19')}（{n.dthr('S19')}）、"
                    f"CoV {n.cov('S19')}、max/min {n.mm('S19')}。",
                    f"**S20**：带宽 {n.thr('S20')}（{n.dthr('S20')}）、"
                    f"CoV {n.cov('S20')}、max/min {n.mm('S20')}。"]),
            dict(t="参考结论",
                 b=["两者都基本保住总带宽，但 CoV 与长期速率差几乎没有离开 S0。"
                    "这说明 requester 窗口能限制注入量，却没有把瓶颈 hop 的服务机会"
                    "从快核搬给慢核。",
                    "[[仅保留为 source-side 控制对照，不进入最终架构建议。]]"])])),

    ("figside", dict(
        chrome="S22 赤字触发的限域让路（13,920 FF‑eq） · 仲裁型 / 总线触发",
        img="24-s22-diagram.png", fig_w=8.75,
        caption="复用 S1 那条 6 bit 广播总线，位宽完全相同，只换了线上放什么："
                "从「我有多堵」换成「我这窗发成了多少」。",
        kicker="RING-ONLY · MECHANISM",
        blocks=[
            dict(kind="card", t="四段机制", b=[
                "**① 检测**：每节点只数自己本窗成功上环的 flit 数，饱和到 6 bit。"
                "播的是进度，不是拥塞等级。",
                "**② 传递**：复用 S1 的 6 bit 广播总线（30 拍）。"
                "自己的进度也从总线读回，两边过同一个量化器、同一段延迟。",
                "**③ 反馈**：赤字 = 总线上 10 项的均值 − 自己那一项，越过 0.5 就举请求。"
                "让路是**单向**的：落后的节点永不让路。",
                "**④ 执行：让位，不是门控**。只对「会从请求者出向 hop 骑过去」的 flit "
                "让路，同时前瞻改发一个会先下环的 flit，让自己的 hop 不空转。"],
                 wt=1.90),
            dict(kind="band",
                 text="S1 用令牌桶：没额度本拍不上环。S22 只在具体某一拍上让出具体某个位置。",
                 wt=0.86),
            dict(kind="band",
                 text="同等不均衡度下让位比门控省 10 倍带宽：S23 门控付 1.50%，"
                      "S22 让位只付 0.15%（K = 2000 同轮对比）。", wt=0.70)])),

    ("media", dict(
        chrome="S22 实测：均匀流量", kicker="MEASURED · UNIFORM",
        stat=n.cov("S22"),
        stat_sub=[f"S22 写侧 100 拍 CoV（S0 = {n.cov('S0')}）",
                  f"总写带宽 {n.dthr('S22')}"],
        img="25-s22-compare.png",
        caption="四根柱依次 S0 / S1 / S16 / S22，橙虚线 = 理论上限 R*；写侧 K = 20000。"
                "后两幅是公平性的两个不同问题，见右侧第三张卡。",
        cards=[
            dict(t="与 S1 同口径", accent=True,
                 b=[f"S22：CoV {n.cov('S22')}，带宽 {n.dthr('S22')}，"
                    f"max/min {n.mm('S22')}，13,920 FF‑eq。",
                    f"S1：CoV {n.cov('S1')}，带宽 {n.dthr('S1')}，"
                    f"max/min {n.mm('S1')}，21,220 FF‑eq。"]),
            dict(t="数字上被 S16 支配，但层次不同",
                 b=[f"S16 用 **1/15 的硬件**拿到 CoV {n.cov('S16')}、"
                    f"带宽 {n.dthr('S16')}。"
                    "S22 唯一的、也是决定性的优势："
                    "[[它只改环上仲裁，一点事务层都不碰。]]"]),
            dict(t="两个公平性指标分别在问什么",
                 b=["**瞬时不均衡度**（每 100 拍算一次十核带宽的 CoV 再聚合）问的是："
                    "在任意一小段时间里，十个核是不是同时都在被服务。",
                    "**长期速率差**（整窗最快核带宽 ÷ 最慢核带宽）问的是："
                    "把整个争用窗平均下来，有没有哪个核被系统性地拖慢。"
                    f"S0 = {n.mm('S0')} 意味着最快的核长期比最慢的核快 "
                    f"{100 * (float(n.mm('S0')) - 1):.0f}%。"])])),

    ("figside", dict(
        chrome="S22 微架构：信号侧 8 个部件 + 仲裁侧插入的两级",
        img="34-s22-uarch.png", fig_w=8.15,
        caption="上 = 信号侧：本窗上环数 → 6 bit 总线 → 10 项表 → 加法树 → 赤字 → 请求 FSM；"
                "下 = 仲裁侧：在已有的 I‑tag 判定之后插入跨越判定与前瞻。",
        kicker="MICROARCHITECTURE · S22",
        blocks=[
            dict(kind="card", t="新增状态（每节点一份）", b=[
                "**ok_win** 6 bit 饱和计数；**cum_bus 表** 10 × 8 bit；"
                "**deficit** ×10，钳位 ±64；",
                "**请求 FSM** ×10：≥ 0.5 举、≤ 0 撤、hold 16 拍到期强制撤。"],
                 wt=1.05),
            dict(kind="card", accent=True, t="仲裁侧插在哪一级", b=[
                "候选 DAT flit 先过已有 I‑tag 判定，再问：有没有请求者 h，"
                "deficit[h] ≥ 我的 + 3.0 且 h 在这个 flit 的剩余路径上？",
                "有 → 往后 8 深找一个目的地不同、不跨任何请求者的 flit 顶替；"
                "找不到才让出本拍。"], wt=1.30),
            dict(kind="band",
                 text="13,920 FF‑eq 里 12,000 是运算，几乎全是那棵 10 输入加法树（360 × 20）。"
                      "请求集合可由每节点持有的同一张表本地重算，不需再加信号。",
                 wt=0.80)])),

    ("figside", dict(
        chrome="S22 时间线：一个 64 拍控制窗里发生了什么",
        img="35-s22-flow.png", fig_w=8.15,
        caption="上 = 事件（时间轴分段缩放，93–109 拍放大）；下 = 落后核 C4 与领先核 C2 "
                "的赤字轨迹（示例数值）。红色阴影 = C4 请求有效期。",
        kicker="WALKTHROUGH · S22",
        blocks=[
            dict(kind="card", t="从测到动要 30 拍", b=[
                "窗末 t = 63 播出 6 bit 计数，t = 93 才送达。这 30 拍里没有人让位 —— "
                "赤字信号永远是 30 拍前的均值，这是它比 S16 慢的结构原因。"],
                 wt=1.05),
            dict(kind="card", accent=True, t="请求怎么撤", b=[
                "上环即扣 1，赤字 ≤ 0 就撤，不等下一窗；",
                "hold 16 拍到期无论如何撤 —— 请求拦不住 transit，"
                "被 transit 卡住的核不能让上游白白空等。"], wt=1.15),
            dict(kind="band",
                 text="margin 3.0：只向比自己至少多落后 3 flit 的请求者让位，"
                      "「差不多齐」的不让，避免白扔 hop。", wt=0.72)])),

    ("figside", dict(
        chrome="S29 日历触发的限域让路（4,440 FF‑eq） · 调度型 / 无拥塞信号",
        img="30-s29-diagram.png", fig_w=8.15,
        caption="上 = 触发：一张 20 拍的固定日历，加每核 1 bit 需求位；"
                "下 = 执行：与 S22 完全同一套环仲裁让位。",
        kicker="ARBITRATION BRANCH · CHEAPEST POINT",
        blocks=[
            dict(kind="card", t="四段机制（与 S1 同格式）", b=[
                "**① 检测**：没有检测。触发是一个 5 bit 帧计数器，常量、"
                "无收敛过程、无控制环路。",
                "**② 传递**：10 bit 需求字（每核 1 bit「我有 WriteData 排队」），"
                "每 16 拍刷新 —— 仅用于回收空闲核的时隙。",
                "**③ 反馈**：轮到某核时，会骑过它出向 hop 的上游核让出一个 slot。",
                "**④ 控制**：让位，不是门控；队列深度维持出厂 8 / 12。"], wt=1.44),
            dict(kind="card", accent=True, t="它是 S22 的廉价版，不是新执行器", b=[
                "执行器与 S22 逐字相同，省下来的是**触发侧**：不播进度、不算均值、"
                "不比赤字，10 项表与 10 输入加法树整块消失 —— "
                "**13,920 → 4,440 FF‑eq**。",
                "代价是它不知道谁落后：日历按核轮转，而不是按「谁被拖慢了」轮转。"],
                 wt=1.16),
            dict(kind="band",
                 text="定位：不在 Pareto 前沿（S16 全面更优），"
                      "但是「只动环仲裁、不碰事务层」这一支上最便宜的点。",
                 wt=0.62)])),

    ("media", dict(
        chrome="S29 实测：均匀流量", kicker="MEASURED · UNIFORM",
        stat=n.cov("S29"),
        stat_sub=[f"S29 写侧 100 拍 CoV（S0 = {n.cov('S0')}）",
                  f"总写带宽 {n.dthr('S29')}"],
        img="31-s29-compare.png",
        caption="五根柱依次 S0 / S1 / S22 / S16 / S29；橙点线 = S1 的水平，"
                "橙虚线 = 理论上限 R*。写侧 K = 20000，与前面各页同口径。",
        cards=[
            dict(t="与 S1 同口径", accent=True,
                 b=[f"S29：带宽 {n.thr('S29')}，CoV {n.cov('S29')}，"
                    f"max/min {n.mm('S29')}，4,440 FF‑eq。",
                    f"S1：带宽 {n.thr('S1')}，CoV {n.cov('S1')}，"
                    f"max/min {n.mm('S1')}，21,220 FF‑eq。"]),
            dict(t="对 S22：便宜 3.1 倍，略差一点",
                 b=[f"CoV {n.cov('S29')} 对 S22 的 {n.cov('S22')}，"
                    f"带宽 {n.dthr('S29')} 对 {n.dthr('S22')}、"
                    f"max/min {n.mm('S29')} 对 {n.mm('S22')}。",
                    "差距的来源就是它省掉的那部分：日历不知道谁落后，"
                    "让位的方向不总是从领先者流向落后者。"]),
            dict(t="非均匀流量下也站得住",
                 b=["hot 上 S29 = R* 的 **98.66%**，与 S22-stock 的 98.75% 相差 "
                    "0.09 个百分点，而硬件只有它的 1/3。",
                    "它不含任何 pattern 先验：日历是常量，需求位是本地事实。"])])),

    ("figside", dict(
        chrome="S29 微架构：删掉 S22 的表与加法树，触发换成日历",
        img="36-s29-uarch.png", fig_w=8.15,
        caption="上 = 触发侧：5 bit 帧计数器查常量表得到持有者，需求位决定该时隙是否有效；"
                "下 = 仲裁侧：与 S22 同一块逻辑，「请求者集合」换成「单一持有者」。",
        kicker="MICROARCHITECTURE · S29",
        blocks=[
            dict(kind="card", t="新增状态（每节点一份）", b=[
                "**帧计数器** 5 bit（t mod 20）；**持有者表** 10 项常量；",
                "**demand 寄存器** 10 bit，30 拍后整字替换；"
                "需求位每 16 拍由「本地 DAT 队列非空」产生。"], wt=1.05),
            dict(kind="card", accent=True, t="仲裁侧逐字复用 S22", b=[
                "跨越判定的第一个操作数从「赤字更大的请求者」换成「本拍持有者」；"
                "不再比赤字，只查一个 bit。",
                "前瞻加深到 32；HA 节点不让位。"], wt=1.15),
            dict(kind="band",
                 text="4,440 FF‑eq = 总线 200 + 计数 240 + 运算 4,000。"
                      "相对 S22 删掉表 1,600 与加法树 / 赤字运算 8,000。", wt=0.72)])),

    ("figside", dict(
        chrome="S29 时间线：日历、需求位与一次让位",
        img="37-s29-flow.png", fig_w=8.15,
        caption="上 = 两帧日历（每核 2 拍）与需求字的采样 / 生效时刻；"
                "下 = 时隙 C6 的两拍里，上游 C2 的两个候选 flit 各自的命运。",
        kicker="WALKTHROUGH · S29",
        blocks=[
            dict(kind="card", t="日历保证什么", b=[
                "每核每 20 拍拿到 2 拍路权，100 拍公平窗里 5 次；硬保证，不靠收敛。",
                "无需求的核其时隙自动作废，那 2 拍谁都可以上。"], wt=1.05),
            dict(kind="card", accent=True, t="需求字是 30–46 拍前的快照", b=[
                "t = 15 采样、t = 45 生效。刚排空的核最多再「占」46 拍时隙；",
                "损失只在上游确有 flit 骑过、且前瞻找不到替代时才发生。"], wt=1.10),
            dict(kind="band",
                 text=f"它不知道谁落后：让位方向按核号轮转，不按赤字。"
                      f"这就是 S29 对 S22 带宽 {n.dthr('S29')} vs {n.dthr('S22')} 的来源。",
                 wt=0.72)])),

    ("section", dict(
        no="08", kicker="SECTION EIGHT", title=["结论与建议"],
        lead="从现象、根因、理论上界到可实现机制，形成完整证据闭环。",
        key_label="BOTTOM LINE",
        key="问题已经收敛为一个架构决策：能改 HA 授权调度就验证 S16；"
            "事务层不能改，则验证只动环上仲裁的 S22。")),

    ("matrix", dict(
        chrome="结论", kicker="CONCLUSIONS", title="三项确定结论 + 一个架构决策",
        cells=[
            dict(t="① 容量问题已闭环", b=[
                f"S0 最忙 hop 利用率 {n.util('S0')}；相对 R* 的 "
                f"{gap_r:.2f}% 缺口已按拍拆成有效载荷、绕环重发和空转，账目无残项。",
                "[[结论：当前问题不是带宽没打满，而是满载时服务机会分配不均。]]"]),
            dict(t="② 公平性是结构问题", b=[
                f"六快四慢随时间稳定存在；即使完全抹平瞬时抖动，CoV 仍有 "
                f"{n.floor_cov('S0')} 的地板，核间长期速率比为 {n.mm('S0')}。",
                "[[结论：只整形发包时机不够，必须在瓶颈处重新分配服务机会。]]"]),
            dict(t="③ 上界与评价口径已建立", b=[
                f"由逐链路占用约束求得 R(CoV) 理论边界与等速率上限 R* = {n.rstar:.4f}；"
                f"边界在 CoV 坐标下是直线，给出单一标量 "
                f"φ = (R − {n.kappa:.3f}·CoV)/R*"
                f"（S16 {n.phi('S16')}、S0 {n.phi('S0')}）；"
                "再用 100 拍 CoV、完成时间比、硬件 FF‑eq 和流量迁移统一评估。",
                "[[结论：方案优劣来自同一上界下的可量化差距，不依赖经验判断。]]"]),
            dict(t="④ 控制点决定方案", accent=True, b=[
                f"S16 在 HA 授权点直接重排服务：900 FF‑eq，CoV {n.cov('S16')}，"
                f"max/min {n.mm('S16')}，总带宽 {n.dthr('S16')}，φ {n.phi('S16')}。",
                f"第 21 页空着的四格已实测：S26 / S27 结构性失效，S28 不动均衡，"
                f"S28S 用带宽 {n.dthr('S28S')} 换均衡；S29 为 {n.triple('S29')}。",
                f"若事务层不可改，只动环仲裁：S22（13,920 FF‑eq，φ {n.phi('S22')}）"
                f"效果最好，S29（4,440，φ {n.phi('S29')}）便宜 3.1 倍；"
                "S19 / S20 仅作对照。",
                "[[架构决策只剩：HA 授权调度能不能改。]]"])],
        band="本研究完成了「现象复现 → 逐拍归因 → 理论上界 → 机制设计 → "
             "硬件代价 → 流量迁移」闭环，并把开放问题压缩为一项可决策的架构边界。")),

    ("process", dict(
        chrome="建议", kicker="RECOMMENDATION",
        title="一主两备，明确研究边界",
        steps=[
            dict(t="主方案 · S16 写侧授权保留", accent=True, b=[
                "**目标**：在 HA / memory controller 内验证 DBIDResp 授权排队与最少服务优先。",
                f"**价值**：当前所有可实现点里，S16 离理论上界最近（φ {n.phi('S16')}）、"
                f"硬件代价最低（900 FF‑eq），uniform 写下带宽 {n.dthr('S16')}。"
                "hot 上 outstanding = 128 会过量注入，先保容量的是窗口类，不是授权重排。",
                "**门槛**：确认 HA 授权时序、CHI 合规性和现有 tracker 接口可支持。"]),
            dict(t="备选甲 · S22 环上赤字让路", b=[
                "**适用条件**：HA 事务调度不可改，但允许调整 ring arbitration。",
                "**验证重点**：6 bit 进度总线复用、30 拍反馈稳定性、"
                "让路范围和现有注入队列深度下的收益。",
                "**定位**：环仲裁分支里效果最好的一档（13,920 FF‑eq）。"]),
            dict(t="备选乙 · S29 日历让路", b=[
                "**适用条件**：同上，且希望进一步压硬件 —— 执行器与 S22 相同，"
                "触发换成 20 拍固定日历 + 每核 1 bit 需求位。",
                f"**代价**：4,440 FF‑eq（S22 的 1/3）；带宽 {n.dthr('S29')} 对 "
                f"S22 的 {n.dthr('S22')}，CoV {n.cov('S29')}。",
                "**定位**：环仲裁分支里硬件最便宜的一档。"]),
            dict(t="边界 · 明确不进入建议", b=[
                "**S19 / S20**：仅作 requester-side 对照；实测公平性几乎不变。",
                "**S26 / S27 / S28**：路由、背压、显式速率三类实测均劣于 S0，"
                "不再推进。",
                f"**读侧**：128 B CompData 下 S0 已达 CoV {n.rcov('S0')}、"
                f"max/min {n.rmm('S0')}，维持现状；"
                f"64 B 读会重现六快四慢（max/min {n.pmm('S0-m1')}），"
                "若 workload 含它则并入 S16-R。"])],
        band="请架构组只决策一件事：是否允许 HA 改变授权的时机与对象。允许则进入 S16 "
             "微架构验证；不允许则转入环仲裁原型，在 S22 与 S29 之间按面积预算取舍。")),

    ("closing", dict(
        title=["汇报完毕", "请架构组决策"],
        lead=["第一问仍是：HA / 内存控制器的授权调度能不能改？",
              "允许修改：进入 S16 微架构验证；不允许修改：进入环仲裁原型"
              "（S22 效果最好，S29 便宜 3.1 倍）。"])),
    ]

BUILDERS = {
    "cover": s_cover, "agenda": s_agenda, "section": s_section,
    "closing": s_closing, "media": s_media, "figside": s_figure_side,
    "process": s_process, "scheme": s_scheme, "matrix": s_matrix,
    "triple": s_triple, "bars": s_bars, "compare": s_compare,
}


def main() -> None:
    if not DECK_JSON.exists():
        raise SystemExit(f"missing {DECK_JSON}; run utils/deck_ring2_data.py")
    n = Live(json.loads(DECK_JSON.read_text()))
    if n.oc != 128:
        raise SystemExit(
            f"{DECK_JSON} has core_outstanding={n.oc}, expected 128")
    deck = slides(n)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    total = len(deck)
    for i, (kind, data) in enumerate(deck, 1):
        slide = BUILDERS[kind](prs, data)
        if kind != "cover":                      # cover footer is a red band
            page_no(slide, i, total, on_dark=(kind == "closing"))
    prs.save(OUT)
    print(f"wrote {OUT}  ({total} slides, outstanding={n.oc})")


if __name__ == "__main__":
    main()
